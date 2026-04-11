import asyncio
import base64
import difflib
import importlib
import imaplib
import ipaddress
import json
import os
import ntpath
import re
import shutil
import socket
import smtplib
import subprocess
import sys
import threading
import time
import traceback
import uuid
from email import message_from_bytes
from email.message import EmailMessage
from email.utils import parseaddr
from collections import deque
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import quote_plus, quote
from urllib.request import urlopen, Request

from main import JarvisLive
from actions.gesture_control import GestureController
from actions.offline_assistant import offline_assistant
from actions.workflow_manager import workflow_manager
from actions.kasa_control import kasa_control
from actions.browser_control import browser_control
from browser_agent_bridge import BrowserAgentBridge
from actions.voice_pipeline import (
    detect_voice_capabilities,
    load_voice_settings,
    save_voice_settings,
    speak_text,
)


BASE_DIR = Path(__file__).resolve().parent
CONFIG_DIR = Path(os.environ.get("BRAHMA_CONFIG_DIR", BASE_DIR / "config"))
CONFIG_DIR.mkdir(parents=True, exist_ok=True)
API_FILE = CONFIG_DIR / "api_keys.json"
ACTION_API_FILE = BASE_DIR / "config" / "api_keys.json"
LEGACY_ACTION_FILE = BASE_DIR / "actions" / "config" / "api_keys.json"
LOG_FILE = os.environ.get("BRAHMA_LOG_FILE")
HOST = "0.0.0.0"
PORT = int(os.environ.get("BRAHMA_BACKEND_PORT", "8770"))
HYBRID_SETTINGS_FILE = CONFIG_DIR / "hybrid_settings.json"
SEQUENCE_SCHEDULE_FILE = CONFIG_DIR / "sequence_schedules.json"
GMAIL_SETTINGS_FILE = CONFIG_DIR / "gmail_settings.json"
DISCORD_SETTINGS_FILE = CONFIG_DIR / "discord_settings.json"
DISCORD_API_BASE = "https://discord.com/api/v10"
PLUGINS_DIR = BASE_DIR / "plugins"

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

if sys.platform == "win32":
    try:
        import ctypes
        import winreg
    except Exception:
        ctypes = None
        winreg = None


class RouteDecision:
    def __init__(self, engine: str, confidence: float, reason: str):
        self.engine = engine
        self.confidence = confidence
        self.reason = reason


def route_command(text: str, advanced_mode: bool = False) -> RouteDecision:
    lowered = str(text or "").lower()
    multi_step_hints = [
        " and ",
        " then ",
        " after ",
        " open ",
        " create ",
        " build ",
        " generate ",
        " download ",
        " search ",
        " website",
        " app ",
        " spreadsheet",
        " presentation",
        " ppt",
        " browser",
        " chrome",
        " gmail",
        " discord",
        " youtube",
    ]
    if any(hint in lowered for hint in multi_step_hints):
        return RouteDecision("brahma_multiaction", 0.62, "Heuristic multi-step routing.")
    return RouteDecision("offline", 0.4, "Default offline assistant routing.")


class PluginManager:
    def __init__(self, root: Path):
        self.root = Path(root)
        self._lock = threading.Lock()
        self._plugins = []
        self._load_errors = []
        self.reload()

    def _load_module(self, entry_path: Path):
        module_name = f"brahma_plugin_{entry_path.stem}_{uuid.uuid4().hex}"
        spec = importlib.util.spec_from_file_location(module_name, entry_path)
        if not spec or not spec.loader:
            raise RuntimeError("Invalid plugin entry spec.")
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)
        return module

    def reload(self):
        with self._lock:
            self._plugins = []
            self._load_errors = []
            if not self.root.exists():
                return
            for plugin_dir in self.root.iterdir():
                if not plugin_dir.is_dir():
                    continue
                manifest_path = plugin_dir / "plugin.json"
                if not manifest_path.exists():
                    continue
                try:
                    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
                except Exception as exc:
                    self._load_errors.append({
                        "name": plugin_dir.name,
                        "error": f"Invalid plugin.json: {exc}",
                    })
                    continue
                entry = str(manifest.get("entry") or "").strip()
                if not entry:
                    self._load_errors.append({
                        "name": manifest.get("name") or plugin_dir.name,
                        "error": "Missing entry in plugin.json.",
                    })
                    continue
                entry_path = (plugin_dir / entry).resolve()
                if entry_path.suffix.lower() != ".py":
                    self._load_errors.append({
                        "name": manifest.get("name") or plugin_dir.name,
                        "error": "Only Python (.py) plugins are supported.",
                    })
                    continue
                if not entry_path.exists():
                    self._load_errors.append({
                        "name": manifest.get("name") or plugin_dir.name,
                        "error": f"Entry not found: {entry}",
                    })
                    continue
                try:
                    module = self._load_module(entry_path)
                    plugin_obj = getattr(module, "plugin", module)
                    handler = getattr(plugin_obj, "on_command", None)
                    if not callable(handler):
                        raise RuntimeError("Plugin entry must export on_command(text, context).")
                    self._plugins.append({
                        "name": manifest.get("name") or plugin_dir.name,
                        "version": manifest.get("version") or "0.0.0",
                        "description": manifest.get("description") or "",
                        "entry": entry,
                        "dir": str(plugin_dir),
                        "handler": handler,
                        "enabled": True,
                        "error": "",
                    })
                except Exception as exc:
                    self._load_errors.append({
                        "name": manifest.get("name") or plugin_dir.name,
                        "error": str(exc),
                    })

    def _run_handler(self, handler, text: str, context: dict):
        result = handler(text, context)
        if asyncio.iscoroutine(result):
            try:
                return asyncio.run(result)
            except RuntimeError:
                loop = asyncio.new_event_loop()
                try:
                    return loop.run_until_complete(result)
                finally:
                    loop.close()
        return result

    def handle(self, text: str, context: dict):
        with self._lock:
            plugins = list(self._plugins)
        for plugin in plugins:
            handler = plugin.get("handler")
            if not callable(handler):
                continue
            try:
                response = self._run_handler(handler, text, context)
                if response:
                    return response, plugin
            except Exception as exc:
                plugin["error"] = str(exc)
        return None, None

    def list_plugins(self):
        with self._lock:
            plugins = [
                {
                    "name": p.get("name"),
                    "version": p.get("version"),
                    "description": p.get("description"),
                    "entry": p.get("entry"),
                    "dir": p.get("dir"),
                    "enabled": p.get("enabled", True),
                    "error": p.get("error", ""),
                }
                for p in self._plugins
            ]
            errors = list(self._load_errors)
        for err in errors:
            plugins.append({
                "name": err.get("name"),
                "version": "0.0.0",
                "description": "",
                "entry": "",
                "dir": "",
                "enabled": False,
                "error": err.get("error") or "Failed to load plugin.",
            })
        return plugins


class BridgeUI:
    def __init__(self):
        self.send_callback = None
        self.logs = deque(maxlen=250)
        self.status_text = "BOOTING"
        self.screen_analysis_active = False
        self.gesture_enabled = False
        self.mic_enabled = False
        self.camera_preview_b64 = None
        self.live_user_text = ""
        self.live_ai_text = ""
        self.saved_sequences = []
        self.kasa_devices = []
        self.voice_settings = self._edge_only_voice_settings(load_voice_settings())
        self.hybrid_settings = self._load_hybrid_settings()
        self.live_ready = False
        self._ensure_primary_api_file()
        self._api_key_ready = self._api_keys_exist()
        self._prime_api_key_env()
        self._lock = threading.Lock()
        self.gesture_controller = GestureController(player=self)
        # Automation / screen control state
        self.automation_mode = "assist"  # observe | assist | do
        self.automation_plan = None
        self.pending_confirmation = False
        self.last_automation_result = None
        self.last_screenshot_b64 = None
        self.advanced_terminal_active = False
        self._last_spoken_message = ""
        self._last_spoken_ts = 0.0
        self.sequence_schedules = self._load_sequence_schedules()
        self.gmail_settings = self._load_gmail_settings()
        self.discord_settings = self._load_discord_settings()
        self.browser_agent_bridge = BrowserAgentBridge(Path(__file__).resolve().parent)
        self.plugin_manager = PluginManager(PLUGINS_DIR)
        self._discord_last_seen = {}
        self._discord_remote_last_seen = {}
        self._discord_remote_warned_channels = set()
        self._discord_bot_user_id = ""
        self._discord_last_command = ""
        self._discord_last_command_ts = 0.0
        self._scheduler_thread = threading.Thread(
            target=self._sequence_scheduler_loop,
            daemon=True,
            name="BrahmaRoutineScheduler",
        )
        self._scheduler_thread.start()
        self._discord_auto_reply_thread = threading.Thread(
            target=self._discord_auto_reply_loop,
            daemon=True,
            name="BrahmaDiscordAutoReply",
        )
        self._discord_auto_reply_thread.start()
        self._discord_remote_thread = threading.Thread(
            target=self._discord_remote_loop,
            daemon=True,
            name="BrahmaDiscordRemote",
        )
        self._discord_remote_thread.start()

    def _edge_only_voice_settings(self, settings=None):
        merged = dict(settings or {})
        # Force runtime speech to Edge TTS for consistent command/readback voice.
        merged["providerOrder"] = ["edge_tts"]
        merged["edgeVoice"] = str(merged.get("edgeVoice") or "en-US-GuyNeural")
        merged["edgeRate"] = str(merged.get("edgeRate") or "-6%")
        merged["edgePitch"] = str(merged.get("edgePitch") or "-2Hz")
        return merged

    def _normalize_schedule_name(self, value: str) -> str:
        return re.sub(r"[^a-z0-9]+", "_", str(value or "").strip().lower()).strip("_")

    def _load_sequence_schedules(self):
        if not SEQUENCE_SCHEDULE_FILE.exists():
            return {}
        try:
            payload = json.loads(SEQUENCE_SCHEDULE_FILE.read_text(encoding="utf-8"))
            if not isinstance(payload, dict):
                return {}
            out = {}
            for key, item in payload.items():
                if not isinstance(item, dict):
                    continue
                name = str(item.get("name") or "").strip()
                run_time = str(item.get("time") or "").strip()
                enabled = bool(item.get("enabled", False))
                last_run_date = str(item.get("last_run_date") or "").strip()
                if name and re.fullmatch(r"\d{2}:\d{2}", run_time):
                    out[self._normalize_schedule_name(key or name)] = {
                        "name": name,
                        "time": run_time,
                        "enabled": enabled,
                        "last_run_date": last_run_date,
                    }
            return out
        except Exception:
            return {}

    def _save_sequence_schedules(self):
        try:
            SEQUENCE_SCHEDULE_FILE.parent.mkdir(parents=True, exist_ok=True)
            SEQUENCE_SCHEDULE_FILE.write_text(
                json.dumps(self.sequence_schedules, indent=2),
                encoding="utf-8",
            )
        except Exception:
            pass

    def _set_sequence_schedule(self, name: str, time_value: str, enabled: bool):
        key = self._normalize_schedule_name(name)
        if not key:
            return
        if enabled and re.fullmatch(r"\d{2}:\d{2}", str(time_value or "")):
            existing = self.sequence_schedules.get(key, {})
            self.sequence_schedules[key] = {
                "name": name,
                "time": str(time_value),
                "enabled": True,
                "last_run_date": str(existing.get("last_run_date") or ""),
            }
        else:
            if key in self.sequence_schedules:
                self.sequence_schedules.pop(key, None)
        self._save_sequence_schedules()

    def _delete_sequence_schedule(self, name: str):
        key = self._normalize_schedule_name(name)
        if key in self.sequence_schedules:
            self.sequence_schedules.pop(key, None)
            self._save_sequence_schedules()

    def _sequence_scheduler_loop(self):
        while True:
            try:
                now = time.localtime()
                hhmm = f"{now.tm_hour:02d}:{now.tm_min:02d}"
                today = time.strftime("%Y-%m-%d", now)
                changed = False
                for key, schedule in list(self.sequence_schedules.items()):
                    if not isinstance(schedule, dict):
                        continue
                    if not schedule.get("enabled"):
                        continue
                    if str(schedule.get("time") or "") != hhmm:
                        continue
                    if str(schedule.get("last_run_date") or "") == today:
                        continue
                    name = str(schedule.get("name") or "").strip()
                    if not name:
                        continue
                    self.write_log(f"[sys] Auto-running routine '{name}' at {hhmm}.")
                    try:
                        result = workflow_manager(
                            {"action": "run", "name": name, "step_runner": self._execute_routine_step},
                            player=self,
                        )
                        self.write_log(f"Brahma AI: {result}")
                    except Exception as exc:
                        self.write_log(f"[error] Auto routine '{name}' failed: {exc}")
                    schedule["last_run_date"] = today
                    changed = True
                if changed:
                    self._save_sequence_schedules()
            except Exception:
                pass
            time.sleep(20)

    def _load_gmail_settings(self):
        settings = {
            "address": str(os.environ.get("BRAHMA_GMAIL_ADDRESS", "")).strip(),
            "app_password": str(os.environ.get("BRAHMA_GMAIL_APP_PASSWORD", "")).strip(),
        }
        try:
            if GMAIL_SETTINGS_FILE.exists():
                payload = json.loads(GMAIL_SETTINGS_FILE.read_text(encoding="utf-8"))
                if isinstance(payload, dict):
                    settings["address"] = str(payload.get("address") or settings["address"]).strip()
                    settings["app_password"] = str(payload.get("app_password") or settings["app_password"]).strip()
        except Exception:
            pass
        return settings

    def _load_discord_settings(self):
        settings = {
            "bot_token": str(os.environ.get("BRAHMA_DISCORD_BOT_TOKEN", "")).strip(),
            "auto_reply_enabled": False,
            "auto_reply_message": "I am offline right now. I will get back to you soon.",
            "auto_reply_channel_ids": [],
            "remote_enabled": False,
            "remote_channel_ids": [],
            "mirror_enabled": True,
        }
        try:
            if DISCORD_SETTINGS_FILE.exists():
                payload = json.loads(DISCORD_SETTINGS_FILE.read_text(encoding="utf-8"))
                if isinstance(payload, dict):
                    settings["bot_token"] = str(payload.get("bot_token") or settings["bot_token"]).strip()
                    settings["auto_reply_enabled"] = bool(payload.get("auto_reply_enabled", False))
                    settings["auto_reply_message"] = str(
                        payload.get("auto_reply_message") or settings["auto_reply_message"]
                    ).strip() or settings["auto_reply_message"]
                    channels = payload.get("auto_reply_channel_ids") or []
                    if isinstance(channels, list):
                        settings["auto_reply_channel_ids"] = [str(item).strip() for item in channels if str(item).strip()]
                    settings["remote_enabled"] = bool(payload.get("remote_enabled", False))
                    remote_channels = payload.get("remote_channel_ids") or []
                    if isinstance(remote_channels, list):
                        settings["remote_channel_ids"] = [str(item).strip() for item in remote_channels if str(item).strip()]
                    settings["mirror_enabled"] = bool(payload.get("mirror_enabled", True))
        except Exception:
            pass
        return settings

    def _save_discord_settings(self):
        try:
            DISCORD_SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
            DISCORD_SETTINGS_FILE.write_text(json.dumps(self.discord_settings, indent=2), encoding="utf-8")
        except Exception:
            pass

    def _gmail_credentials(self):
        self.gmail_settings = self._load_gmail_settings()
        address = str(self.gmail_settings.get("address") or "").strip()
        app_password = str(self.gmail_settings.get("app_password") or "").strip()
        return address, app_password

    def _discord_token(self) -> str:
        latest = self._load_discord_settings()
        if latest.get("bot_token"):
            self.discord_settings["bot_token"] = latest.get("bot_token")
        return str(self.discord_settings.get("bot_token") or "").strip()

    def _discord_api_request(self, method: str, path: str, payload=None, query=None, timeout: int = 20):
        token = self._discord_token()
        if not token:
            raise RuntimeError(
                "Discord bot token is missing. Set BRAHMA_DISCORD_BOT_TOKEN or add bot_token in config/discord_settings.json."
            )
        clean_path = "/" + str(path or "").lstrip("/")
        url = f"{DISCORD_API_BASE}{clean_path}"
        if isinstance(query, dict) and query:
            parts = []
            for key, value in query.items():
                if value is None:
                    continue
                parts.append(f"{quote_plus(str(key))}={quote_plus(str(value))}")
            if parts:
                url = f"{url}?{'&'.join(parts)}"
        body = None
        headers = {
            "Authorization": f"Bot {token}",
            "User-Agent": "BrahmaAI/4.0",
            "Content-Type": "application/json",
        }
        if payload is not None:
            body = json.dumps(payload).encode("utf-8")
        req = Request(url, data=body, method=method.upper(), headers=headers)
        try:
            with urlopen(req, timeout=timeout) as response:
                raw = response.read().decode("utf-8", errors="replace")
                if not raw.strip():
                    return {}
                return json.loads(raw)
        except Exception as exc:
            detail = ""
            try:
                if hasattr(exc, "read"):
                    detail = exc.read().decode("utf-8", errors="replace")
            except Exception:
                detail = ""
            message = str(exc)
            if detail:
                message = f"{message} | {detail[:260]}"
            raise RuntimeError(message) from exc

    def _discord_get_bot_user(self):
        payload = self._discord_api_request("GET", "/users/@me")
        user_id = str(payload.get("id") or "").strip()
        if user_id:
            self._discord_bot_user_id = user_id
        return payload

    def _discord_list_guilds(self):
        payload = self._discord_api_request("GET", "/users/@me/guilds", query={"limit": 200})
        if isinstance(payload, list):
            return payload
        return []

    def _discord_resolve_guild(self, name_or_id: str):
        query = str(name_or_id or "").strip()
        if not query:
            return None
        guilds = self._discord_list_guilds()
        if query.isdigit():
            for guild in guilds:
                if str(guild.get("id")) == query:
                    return guild
        normalized_query = re.sub(r"[^a-z0-9]+", " ", query.lower()).strip()
        best = None
        best_score = 0.0
        for guild in guilds:
            name = str(guild.get("name") or "")
            normalized_name = re.sub(r"[^a-z0-9]+", " ", name.lower()).strip()
            if not normalized_name:
                continue
            score = difflib.SequenceMatcher(None, normalized_query, normalized_name).ratio()
            if normalized_query and normalized_query in normalized_name:
                score += 0.35
            if score > best_score:
                best = guild
                best_score = score
        return best if best_score >= 0.45 else None

    def _discord_list_text_channels(self, guild_id: str):
        payload = self._discord_api_request("GET", f"/guilds/{guild_id}/channels")
        if not isinstance(payload, list):
            return []
        return [item for item in payload if int(item.get("type", -1)) == 0]

    def _discord_resolve_channel(self, guild_id: str, channel_hint: str = ""):
        channels = self._discord_list_text_channels(guild_id)
        if not channels:
            return None
        hint = str(channel_hint or "").strip().lstrip("#")
        if not hint:
            return channels[0]
        normalized_hint = re.sub(r"[^a-z0-9]+", " ", hint.lower()).strip()
        best = None
        best_score = 0.0
        for channel in channels:
            name = str(channel.get("name") or "")
            normalized_name = re.sub(r"[^a-z0-9]+", " ", name.lower()).strip()
            if not normalized_name:
                continue
            score = difflib.SequenceMatcher(None, normalized_hint, normalized_name).ratio()
            if normalized_hint and normalized_hint in normalized_name:
                score += 0.35
            if score > best_score:
                best = channel
                best_score = score
        return best if best and best_score >= 0.4 else channels[0]

    def _discord_send_channel_message(self, channel_id: str, content: str):
        text = str(content or "").strip()
        if not text:
            raise RuntimeError("Message text is required.")
        self._discord_api_request("POST", f"/channels/{channel_id}/messages", payload={"content": text})

    def _discord_send_dm(self, user_id: str, content: str):
        if not str(user_id or "").strip().isdigit():
            raise RuntimeError("DM requires a numeric Discord user id.")
        dm = self._discord_api_request("POST", "/users/@me/channels", payload={"recipient_id": str(user_id).strip()})
        channel_id = str(dm.get("id") or "").strip()
        if not channel_id:
            raise RuntimeError("Could not open DM channel for that user.")
        self._discord_send_channel_message(channel_id, content)
        return channel_id

    def _discord_extract_message_body(self, text: str):
        raw = str(text or "")
        quoted = re.search(r"[\"']([^\"']{1,1800})[\"']", raw)
        if quoted:
            return quoted.group(1).strip()
        patterns = [
            r"\b(?:message|say|text|send)\b\s*(?:to\s+[^,]+)?\s*(?:in\s+[^,]+)?\s*(.+)$",
            r"\bwith\b\s+(.+)$",
        ]
        for pattern in patterns:
            match = re.search(pattern, raw, flags=re.IGNORECASE)
            if match:
                body = match.group(1).strip(" .")
                if body:
                    return body
        return ""

    def _discord_extract_server_name(self, text: str):
        raw = str(text or "")
        for pattern in [
            r"\bserver\s*(?:named|called)?\s*[\"']([^\"']+)['\"]",
            r"\bin\s+server\s+([a-zA-Z0-9 _-]{2,80})",
            r"\bserver\s+([a-zA-Z0-9 _-]{2,80})",
        ]:
            match = re.search(pattern, raw, flags=re.IGNORECASE)
            if match:
                return match.group(1).strip(" .")
        return ""

    def _discord_extract_channel_name(self, text: str):
        raw = str(text or "")
        hashtag = re.search(r"#([a-zA-Z0-9_-]{1,80})", raw)
        if hashtag:
            return hashtag.group(1).strip()
        channel = re.search(r"\bchannel\s+([a-zA-Z0-9 _-]{1,80})", raw, flags=re.IGNORECASE)
        if channel:
            return channel.group(1).strip()
        return ""

    def _discord_auto_reply_loop(self):
        while True:
            try:
                settings = dict(self.discord_settings or {})
                if not settings.get("auto_reply_enabled"):
                    time.sleep(12)
                    continue
                token = self._discord_token()
                auto_message = str(settings.get("auto_reply_message") or "").strip()
                if not token or not auto_message:
                    time.sleep(12)
                    continue
                bot = self._discord_get_bot_user()
                bot_id = str(bot.get("id") or "").strip()
                bot_mention = f"<@{bot_id}>" if bot_id else ""

                dm_channels_payload = self._discord_api_request("GET", "/users/@me/channels")
                dm_channels = {}
                if isinstance(dm_channels_payload, list):
                    for channel in dm_channels_payload:
                        if int(channel.get("type", -1)) == 1:
                            channel_id = str(channel.get("id") or "").strip()
                            if channel_id:
                                dm_channels[channel_id] = True

                monitored = set(str(item).strip() for item in settings.get("auto_reply_channel_ids") or [] if str(item).strip())
                monitored.update(dm_channels.keys())
                for channel_id in monitored:
                    messages = self._discord_api_request("GET", f"/channels/{channel_id}/messages", query={"limit": 8})
                    if not isinstance(messages, list):
                        continue
                    last_seen = str(self._discord_last_seen.get(channel_id) or "")
                    if not last_seen and messages:
                        first_id = str((messages[0] or {}).get("id") or "").strip()
                        if first_id:
                            self._discord_last_seen[channel_id] = first_id
                        continue
                    newest_seen = last_seen
                    for msg in reversed(messages):
                        msg_id = str(msg.get("id") or "")
                        if not msg_id:
                            continue
                        if newest_seen and msg_id <= newest_seen:
                            continue
                        author = msg.get("author") or {}
                        author_id = str(author.get("id") or "")
                        if not author_id or author_id == bot_id or bool(author.get("bot")):
                            newest_seen = msg_id
                            continue
                        content = str(msg.get("content") or "")
                        is_dm = channel_id in dm_channels
                        if is_dm or (bot_mention and bot_mention in content):
                            self._discord_send_channel_message(channel_id, auto_message)
                        newest_seen = msg_id
                    if newest_seen and newest_seen != last_seen:
                        self._discord_last_seen[channel_id] = newest_seen
            except Exception:
                pass
            time.sleep(18)

    def _discord_remote_loop(self):
        while True:
            try:
                settings = dict(self.discord_settings or {})
                if not settings.get("remote_enabled"):
                    time.sleep(10)
                    continue
                token = self._discord_token()
                if not token:
                    time.sleep(12)
                    continue
                channel_ids = [
                    str(item).strip()
                    for item in settings.get("remote_channel_ids") or []
                    if str(item).strip()
                ]
                if not channel_ids:
                    time.sleep(12)
                    continue
                bot = self._discord_get_bot_user()
                bot_id = str(bot.get("id") or "").strip()
                for channel_id in channel_ids:
                    messages = self._discord_api_request(
                        "GET", f"/channels/{channel_id}/messages", query={"limit": 8}
                    )
                    if not isinstance(messages, list):
                        continue
                    last_seen = str(self._discord_remote_last_seen.get(channel_id) or "")
                    now_ms = int(time.time() * 1000)
                    newest_seen = last_seen
                    for msg in reversed(messages):
                        msg_id = str(msg.get("id") or "")
                        if not msg_id:
                            continue
                        if newest_seen and int(msg_id) <= int(newest_seen):
                            continue
                        author = msg.get("author") or {}
                        author_id = str(author.get("id") or "")
                        if not author_id or author_id == bot_id or bool(author.get("bot")):
                            newest_seen = msg_id
                            continue
                        content = str(msg.get("content") or "").strip()
                        if not content:
                            if channel_id not in self._discord_remote_warned_channels:
                                self._discord_remote_warned_channels.add(channel_id)
                                self._discord_send_channel_message(
                                    channel_id,
                                    "Discord remote is online, but I can't read message content. Enable Message Content Intent in the bot settings.",
                                )
                            newest_seen = msg_id
                            continue
                        msg_ts = self._discord_snowflake_time_ms(msg_id)
                        if not last_seen and msg_ts and (now_ms - msg_ts) > 120000:
                            newest_seen = msg_id
                            continue
                        response = self.process_command_text(
                            content,
                            source="discord",
                            discord_channel_id=channel_id,
                        )
                        if response:
                            self._discord_send_channel_message(channel_id, response)
                        newest_seen = msg_id
                    if newest_seen and newest_seen != last_seen:
                        self._discord_remote_last_seen[channel_id] = newest_seen
            except Exception:
                pass
            time.sleep(6)

    def _discord_snowflake_time_ms(self, value: str) -> int:
        try:
            snowflake = int(str(value).strip())
        except Exception:
            return 0
        return (snowflake >> 22) + 1420070400000

    def _load_hybrid_settings(self):
        defaults = {
            "advanced_mode_enabled": False,
            "project_workspace_path": "",
            "project_workspace_name": "",
        }
        if not HYBRID_SETTINGS_FILE.exists():
            return defaults
        try:
            with open(HYBRID_SETTINGS_FILE, "r", encoding="utf-8") as handle:
                payload = json.load(handle)
            if isinstance(payload, dict):
                merged = self._sanitize_hybrid_settings({**defaults, **payload})
                if merged != {**defaults, **payload}:
                    HYBRID_SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
                    with open(HYBRID_SETTINGS_FILE, "w", encoding="utf-8") as handle:
                        json.dump(merged, handle, indent=2)
                return merged
        except Exception:
            pass
        return defaults

    def _sanitize_hybrid_settings(self, payload):
        merged = dict(payload or {})
        project_path = str(merged.get("project_workspace_path") or "").strip()
        if project_path:
            project_target = Path(project_path).expanduser()
            if not project_target.exists() or not project_target.is_dir():
                merged["project_workspace_path"] = ""
                merged["project_workspace_name"] = ""
        else:
            merged["project_workspace_path"] = ""
            merged["project_workspace_name"] = ""
        return merged

    def save_hybrid_settings(self, next_settings):
        filtered = {
            key: value
            for key, value in (next_settings or {}).items()
            if value is not None
        }
        merged = {**self.hybrid_settings, **filtered}
        merged["advanced_mode_enabled"] = bool(merged.get("advanced_mode_enabled", False))
        merged["project_workspace_path"] = str(merged.get("project_workspace_path") or "").strip()
        merged["project_workspace_name"] = str(merged.get("project_workspace_name") or "").strip()
        merged = self._sanitize_hybrid_settings(merged)
        HYBRID_SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
        with open(HYBRID_SETTINGS_FILE, "w", encoding="utf-8") as handle:
            json.dump(merged, handle, indent=2)
        self.hybrid_settings = merged
        return self.hybrid_settings

    def get_hybrid_state(self):
        return {
            "advancedMode": bool(self.hybrid_settings.get("advanced_mode_enabled", False)),
            "projectWorkspacePath": self.hybrid_settings.get("project_workspace_path"),
            "projectWorkspaceName": self.hybrid_settings.get("project_workspace_name"),
            "openclawEnabled": False,
            "openclawStatus": {"state": "removed"},
        }

    def _inject_project_context(self, text: str) -> str:
        project_path = str(self.hybrid_settings.get("project_workspace_path") or "").strip()
        if not project_path:
            return text
        project_name = str(self.hybrid_settings.get("project_workspace_name") or "").strip() or Path(project_path).name
        return (
            f"Current Brahma project workspace: {project_path} ({project_name}). "
            "For any code, websites, apps, scaffolding, files, assets, folders, downloads, or generated deliverables, "
            "work inside this project workspace unless I explicitly ask for another location.\n\n"
            f"User request: {text}"
        )

    def _extract_first_url(self, text: str) -> str:
        match = re.search(r"https?://[^\s\"'<>]+", str(text or ""))
        return match.group(0).strip() if match else ""

    def _open_external_url(self, url: str) -> bool:
        target = str(url or "").strip()
        if not target:
            return False
        try:
            if sys.platform == "win32":
                subprocess.Popen(["cmd", "/c", "start", "", target], shell=False)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", target])
            else:
                subprocess.Popen(["xdg-open", target])
            return True
        except Exception:
            return False

    def _open_in_managed_browser(self, url: str, allow_external_fallback: bool = True):
        target = str(url or "").strip()
        if not target:
            return False, "Empty URL."
        browser_result = ""
        try:
            browser_result = str(
                browser_control(
                    parameters={"action": "go_to", "url": target},
                    player=self,
                )
                or ""
            ).strip()
            lower = browser_result.lower()
            hard_failure = any(
                token in lower
                for token in (
                    "unknown action",
                    "navigation error",
                    "not started",
                    "traceback",
                )
            )
            if not hard_failure:
                return True, browser_result or f"Opened: {target}"
        except Exception as exc:
            browser_result = str(exc or "").strip()

        if allow_external_fallback and self._open_external_url(target):
            return True, f"Opened externally: {target}"
        return False, browser_result or "Could not open URL."

    def _open_external_path(self, path_value: str) -> bool:
        target = str(path_value or "").strip()
        if not target:
            return False
        try:
            if sys.platform == "win32":
                subprocess.Popen(["explorer", target], shell=False)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", target])
            else:
                subprocess.Popen(["xdg-open", target])
            return True
        except Exception:
            return False

    def _current_project_workspace(self) -> Path | None:
        project_path = str(self.hybrid_settings.get("project_workspace_path") or "").strip()
        if not project_path:
            return None
        try:
            target = Path(project_path).expanduser()
            target.mkdir(parents=True, exist_ok=True)
            if target.is_dir():
                return target
        except Exception:
            return None
        return None

    def _extract_named_project(self, text: str) -> str:
        raw = str(text or "")
        quoted = re.search(r"[\"']([^\"']{2,80})[\"']", raw)
        if quoted:
            return quoted.group(1).strip()
        named = re.search(r"\b(?:called|named|name\s+it)\s+([a-zA-Z0-9 _-]{2,80})", raw, flags=re.IGNORECASE)
        if named:
            return named.group(1).strip()
        return ""

    def _extract_brand_name(self, text: str) -> str:
        raw = str(text or "")
        named = re.search(
            r"\b(?:name|called|brand|website name)\s*(?:is|:)?\s*([a-zA-Z0-9 _-]{2,60})",
            raw,
            flags=re.IGNORECASE,
        )
        if named:
            return named.group(1).strip()
        return ""

    def _extract_topic(self, text: str) -> str:
        raw = str(text or "")
        patterns = [
            r"\btopic\s*[-:]\s*[\"']?([^\"'\n]+)",
            r"\bon topic\b\s*[\"']?([^\"'\n]+)",
            r"\babout\b\s+[\"']?([^\"'\n]+)",
            r"\bon\s+([a-zA-Z0-9][a-zA-Z0-9 _-]{1,120}?)(?:\s+in\s+(?:downloads?|desktop|documents?|project\s+(?:workspace|folder))\b|$)",
            r"\bfor\s+([a-zA-Z0-9][a-zA-Z0-9 _-]{1,120}?)(?:\s+in\s+(?:downloads?|desktop|documents?|project\s+(?:workspace|folder))\b|$)",
        ]
        for pattern in patterns:
            hit = re.search(pattern, raw, flags=re.IGNORECASE)
            if hit:
                topic = hit.group(1).strip(" .,:;\"'")
                topic = re.sub(
                    r"\b(?:wallpaper|image|images|photo|photos|download|downloads|folder|set|it|as)\b",
                    " ",
                    topic,
                    flags=re.IGNORECASE,
                )
                topic = re.sub(r"\s+", " ", topic).strip(" .,:;\"'")
                if topic:
                    return topic
        media_hint = re.search(
            r"\b([a-zA-Z0-9][a-zA-Z0-9 _-]{1,80})\s+(?:wallpaper|images?|photos?)\b",
            raw,
            flags=re.IGNORECASE,
        )
        if media_hint:
            topic = media_hint.group(1).strip(" .,:;\"'")
            topic = re.sub(
                r"\b(?:download|me|a|an|the|good|best|high quality|hd|4k|get|please)\b",
                " ",
                topic,
                flags=re.IGNORECASE,
            )
            topic = re.sub(r"\s+", " ", topic).strip(" .,:;\"'")
            if topic:
                return topic
        return "AI"

    def _friendly_builder_reason(self, reason: str) -> str:
        msg = str(reason or "").strip()
        lower = msg.lower()
        if any(token in lower for token in ("quota", "rate limit", "resourceexhausted", "429", "exceeded your current quota")):
            return "AI quota limit reached, so fallback mode was used."
        if "service is currently unavailable" in lower:
            return "AI service was temporarily unavailable, so fallback mode was used."
        if not msg:
            return "AI planner unavailable; fallback mode was used."
        compact = re.sub(r"\s+", " ", msg).strip()
        if len(compact) > 140:
            compact = compact[:140].rstrip() + "..."
        return f"Fallback mode was used. ({compact})"

    def _safe_path_name(self, value: str, fallback: str = "Brahma Output") -> str:
        cleaned = re.sub(r'[<>:"/\\|?*]+', " ", str(value or ""))
        cleaned = re.sub(r"\s+", " ", cleaned).strip(" .")
        return cleaned[:80] if cleaned else fallback

    def _extract_folder_name(self, text: str) -> str:
        raw = str(text or "")
        patterns = [
            r"\bfolder\s+named\s+[\"']([^\"']+)[\"']",
            r"\bfolder\s+named\s+([a-zA-Z0-9 _-]{2,80})",
            r"\bin\s+folder\s+[\"']([^\"']+)[\"']",
            r"\bin\s+([a-zA-Z0-9 _-]{2,80})\s+folder\b",
        ]
        for pattern in patterns:
            hit = re.search(pattern, raw, flags=re.IGNORECASE)
            if hit:
                name = str(hit.group(1) or "").strip()
                name = re.sub(
                    r"\s+in\s+(downloads?|desktop|documents?|project\s+(?:workspace|folder))\b.*$",
                    "",
                    name,
                    flags=re.IGNORECASE,
                ).strip()
                if re.fullmatch(r"(downloads?|desktop|documents?|project\s+(?:workspace|folder))", name, flags=re.IGNORECASE):
                    return ""
                return self._safe_path_name(name)
        return ""

    def _extract_spreadsheet_topic(self, text: str) -> str:
        raw = str(text or "")
        quoted = re.search(r"[\"']([^\"']{2,120})[\"']", raw)
        if quoted:
            return quoted.group(1).strip()
        patterns = [
            r"\btopic\s*[-:]\s*([a-zA-Z0-9 _-]{2,120})",
            r"\bon topic\b\s*([a-zA-Z0-9 _-]{2,120})",
            r"\babout\b\s*([a-zA-Z0-9 _-]{2,120})",
        ]
        for pattern in patterns:
            hit = re.search(pattern, raw, flags=re.IGNORECASE)
            if hit:
                candidate = hit.group(1).strip(" .,:;")
                if candidate:
                    return candidate
        fallback = self._extract_topic(raw)
        return fallback or "AI"

    def _resolve_export_root(self, text: str) -> Path:
        lower = str(text or "").lower()
        home = Path.home()
        if "project workspace" in lower or "project folder" in lower:
            workspace = self._current_project_workspace()
            if workspace:
                return workspace
        if "desktop" in lower:
            return home / "Desktop"
        if "document" in lower:
            return home / "Documents"
        return home / "Downloads"

    def _resolve_export_folder(self, text: str, default_name: str) -> Path:
        root = self._resolve_export_root(text)
        root.mkdir(parents=True, exist_ok=True)
        folder_name = self._extract_folder_name(text) or default_name
        target = root / self._safe_path_name(folder_name, fallback=default_name)
        target.mkdir(parents=True, exist_ok=True)
        return target

    def _fetch_json(self, url: str):
        req = Request(url, headers={"User-Agent": "BrahmaAI/4.0"})
        with urlopen(req, timeout=15) as response:
            return json.loads(response.read().decode("utf-8", errors="replace"))

    def _research_topic(self, topic: str):
        topic_clean = str(topic or "").strip() or "AI"
        summary = ""
        source_url = ""
        points = []
        rows = []
        errors = []
        try:
            summary_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{quote(topic_clean)}"
            payload = self._fetch_json(summary_url)
            summary = str(payload.get("extract") or "").strip()
            source_url = str((((payload.get("content_urls") or {}).get("desktop") or {}).get("page")) or "")
        except Exception as exc:
            errors.append(f"summary:{exc}")

        try:
            search_url = (
                "https://en.wikipedia.org/w/api.php?action=query&list=search"
                f"&srsearch={quote_plus(topic_clean)}&utf8=&format=json&srlimit=12"
            )
            payload = self._fetch_json(search_url)
            for item in (payload.get("query") or {}).get("search", []):
                title = str(item.get("title") or "").strip()
                snippet = re.sub(r"<[^>]+>", "", str(item.get("snippet") or "")).strip()
                if not title:
                    continue
                rows.append({
                    "title": title,
                    "snippet": snippet,
                    "url": f"https://en.wikipedia.org/wiki/{quote(title.replace(' ', '_'))}",
                })
        except Exception as exc:
            errors.append(f"search:{exc}")

        if summary:
            points = [s.strip() for s in re.split(r"(?<=[.!?])\s+", summary) if s.strip()]
            points = points[:14]
        if not points and rows:
            points = [f"{row['title']}: {row['snippet']}" for row in rows[:10] if row.get("snippet")]
        if not summary and rows:
            summary = " ".join(filter(None, [rows[0].get("snippet", ""), rows[1].get("snippet", "")])).strip()
        if not source_url and rows:
            source_url = rows[0].get("url", "")

        if not summary:
            summary = (
                f"Automated web research was limited for topic '{topic_clean}'. "
                "A starter spreadsheet was still generated and is ready for manual enrichment."
            )
        if not points:
            points = [f"{topic_clean} overview placeholder", f"Add custom research notes for {topic_clean}"]

        return {
            "topic": topic_clean,
            "summary": summary,
            "source_url": source_url,
            "points": points,
            "rows": rows[:30],
            "research_limited": bool(errors),
        }

    def _ensure_openpyxl(self):
        try:
            import openpyxl  # type: ignore
            return openpyxl
        except Exception:
            pass
        try:
            subprocess.run(
                [sys.executable, "-m", "pip", "install", "--disable-pip-version-check", "-q", "openpyxl"],
                capture_output=True,
                text=True,
                timeout=180,
            )
            import openpyxl  # type: ignore
            return openpyxl
        except Exception:
            return None

    def _write_csv(self, path_obj: Path, headers, rows):
        import csv
        path_obj.parent.mkdir(parents=True, exist_ok=True)
        with open(path_obj, "w", newline="", encoding="utf-8") as handle:
            writer = csv.writer(handle)
            writer.writerow(headers)
            for row in rows:
                writer.writerow(row)

    def _save_research_spreadsheets(self, folder: Path, research: dict):
        topic_slug = self._slugify_name(research.get("topic", "topic")) or "topic"
        now_text = time.strftime("%Y-%m-%d %H:%M:%S")
        files = []
        openpyxl = self._ensure_openpyxl()
        if openpyxl:
            def _save_book(path_obj: Path, headers, rows):
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.title = "Data"
                ws.append(headers)
                for row in rows:
                    ws.append(list(row))
                wb.save(path_obj)
            summary_path = folder / f"{topic_slug}_summary.xlsx"
            _save_book(
                summary_path,
                ["Topic", "Summary", "Primary Source", "Generated At"],
                [[research["topic"], research["summary"], research.get("source_url", ""), now_text]],
            )
            points_path = folder / f"{topic_slug}_key_points.xlsx"
            _save_book(
                points_path,
                ["#", "Key Point"],
                [[idx + 1, point] for idx, point in enumerate(research.get("points", []))],
            )
            rows_path = folder / f"{topic_slug}_research_results.xlsx"
            _save_book(
                rows_path,
                ["Title", "Snippet", "URL"],
                [[row.get("title", ""), row.get("snippet", ""), row.get("url", "")] for row in research.get("rows", [])],
            )
            files = [summary_path, points_path, rows_path]
            return files, "xlsx"

        summary_csv = folder / f"{topic_slug}_summary.csv"
        points_csv = folder / f"{topic_slug}_key_points.csv"
        rows_csv = folder / f"{topic_slug}_research_results.csv"
        self._write_csv(
            summary_csv,
            ["Topic", "Summary", "Primary Source", "Generated At"],
            [[research["topic"], research["summary"], research.get("source_url", ""), now_text]],
        )
        self._write_csv(
            points_csv,
            ["#", "Key Point"],
            [[idx + 1, point] for idx, point in enumerate(research.get("points", []))],
        )
        self._write_csv(
            rows_csv,
            ["Title", "Snippet", "URL"],
            [[row.get("title", ""), row.get("snippet", ""), row.get("url", "")] for row in research.get("rows", [])],
        )
        files = [summary_csv, points_csv, rows_csv]
        return files, "csv"

    def _handle_spreadsheet_request(self, text: str) -> str:
        topic = self._extract_spreadsheet_topic(text)
        folder = self._resolve_export_folder(text, default_name=f"{topic} Research")
        research = self._research_topic(topic)
        files, fmt = self._save_research_spreadsheets(folder, research)
        self._open_external_path(str(folder))
        limited_note = " (research fallback used)" if research.get("research_limited") else ""
        file_list = ", ".join(path.name for path in files)
        return (
            f"Created {len(files)} {fmt.upper()} research spreadsheets for '{topic}' in {folder}{limited_note}. "
            f"Files: {file_list}"
        )

    def _ensure_python_package(self, package_name: str, import_name: str = "") -> bool:
        module_name = import_name or package_name.replace("-", "_")
        try:
            importlib.import_module(module_name)
            return True
        except Exception:
            pass
        try:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "--disable-pip-version-check", "-q", package_name],
                capture_output=True,
                text=True,
                timeout=240,
            )
            if result.returncode != 0:
                return False
            importlib.import_module(module_name)
            return True
        except Exception:
            return False

    def _extract_slide_count(self, text: str, default_count: int = 10) -> int:
        raw = str(text or "")
        match = re.search(r"\b(\d{1,2})\s*(?:slides?|pages?)\b", raw, flags=re.IGNORECASE)
        if match:
            try:
                value = int(match.group(1))
                return max(1, min(25, value))
            except Exception:
                pass
        return default_count

    def _download_topic_images(self, topic: str, target_dir: Path, count: int = 5):
        target_dir.mkdir(parents=True, exist_ok=True)
        saved_paths = []
        normalized_topic = str(topic or "").strip() or "topic"
        safe_topic = self._safe_path_name(normalized_topic, fallback="topic").replace(" ", "_").lower()
        query = re.sub(
            r"\b(wallpaper|images?|photos?|picture|pics?|download|set|apply|background)\b",
            " ",
            normalized_topic,
            flags=re.IGNORECASE,
        )
        query = re.sub(r"\s+", " ", query).strip() or normalized_topic
        query_lower = query.lower()
        avoid_people = not any(
            token in query_lower for token in ("person", "people", "portrait", "face", "girl", "boy", "woman", "man", "character")
        )
        banned_terms = {"girl", "woman", "women", "model", "selfie", "portrait", "face", "actress", "celebrity"}

        def _download(url: str, index: int):
            req = Request(url, headers={"User-Agent": "BrahmaAI/4.0"})
            with urlopen(req, timeout=18) as response:
                content_type = str(response.headers.get("Content-Type") or "").lower()
                ext = ".jpg"
                if "png" in content_type:
                    ext = ".png"
                elif "webp" in content_type:
                    ext = ".webp"
                path_obj = target_dir / f"{safe_topic}_{index:02d}{ext}"
                path_obj.write_bytes(response.read())
                return path_obj

        commons_url = (
            "https://commons.wikimedia.org/w/api.php?action=query&generator=search"
            f"&gsrsearch={quote_plus(query + ' wallpaper filetype:bitmap')}&gsrnamespace=6&gsrlimit={max(8, count * 3)}"
            "&prop=imageinfo&iiprop=url&format=json"
        )
        try:
            payload = self._fetch_json(commons_url)
            pages = (payload.get("query") or {}).get("pages") or {}
            urls = []
            for page in pages.values():
                title = str((page or {}).get("title") or "").lower()
                if avoid_people and any(term in title for term in banned_terms):
                    continue
                image_info = (page or {}).get("imageinfo") or []
                if image_info and image_info[0].get("url"):
                    urls.append(str(image_info[0]["url"]))
            for idx, image_url in enumerate(urls[:count], start=1):
                try:
                    saved_paths.append(_download(image_url, idx))
                except Exception:
                    continue
        except Exception:
            pass

        if len(saved_paths) < count:
            try:
                openverse_query = f"{query} wallpaper"
                api_url = (
                    "https://api.openverse.org/v1/images/"
                    f"?q={quote_plus(openverse_query)}&page_size={max(20, count * 8)}"
                    "&license=cc0%2Cby%2Cby-sa%2Cpdm"
                )
                payload = self._fetch_json(api_url)
                results = payload.get("results") or []
                ranked = []
                for item in results:
                    title = str(item.get("title") or "").lower()
                    creator = str(item.get("creator") or "").lower()
                    if avoid_people and any(term in title or term in creator for term in banned_terms):
                        continue
                    width = int(item.get("width") or 0)
                    height = int(item.get("height") or 0)
                    score = width * height
                    if "wallpaper" in title:
                        score += 200000
                    if "game" in title and "gaming" in query_lower:
                        score += 250000
                    if width >= 1280:
                        score += 100000
                    src = str(item.get("url") or item.get("source") or "").strip()
                    if src.startswith("http"):
                        ranked.append((score, src))
                ranked.sort(key=lambda x: x[0], reverse=True)
                next_index = len(saved_paths) + 1
                for _, image_url in ranked:
                    if len(saved_paths) >= count:
                        break
                    try:
                        saved_paths.append(_download(image_url, next_index))
                        next_index += 1
                    except Exception:
                        continue
            except Exception:
                pass
        return saved_paths

    def _extract_ppt_topic(self, text: str) -> str:
        raw = str(text or "")
        quoted = re.search(r"[\"']([^\"']{2,120})[\"']", raw)
        if quoted:
            return quoted.group(1).strip()
        for pattern in (
            r"\btopic\s*[-:]\s*([a-zA-Z0-9 _-]{2,120})",
            r"\bon\s+topic\b\s*([a-zA-Z0-9 _-]{2,120})",
            r"\babout\b\s*([a-zA-Z0-9 _-]{2,120})",
            r"\bon\b\s*([a-zA-Z0-9 _-]{2,120})",
        ):
            hit = re.search(pattern, raw, flags=re.IGNORECASE)
            if hit:
                return str(hit.group(1) or "").strip(" .,:;")
        return "AI"

    def _handle_ppt_request(self, text: str) -> str:
        topic = self._extract_ppt_topic(text)
        slide_count = self._extract_slide_count(text, default_count=10)
        folder = self._resolve_export_folder(text, default_name=f"{topic} Presentation")
        if not self._ensure_python_package("python-pptx", "pptx"):
            return (
                "PowerPoint generation needs `python-pptx`, but I couldn't install it automatically. "
                "Run once: python -m pip install python-pptx"
            )
        self._ensure_python_package("Pillow", "PIL")

        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore

        research = self._research_topic(topic)
        points = list(research.get("points") or [])
        if not points:
            points = [f"{topic} overview", f"{topic} use cases", f"Future of {topic}"]
        images = self._download_topic_images(topic, folder / "images", count=min(6, max(1, slide_count - 1)))

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Designed by Brahma AI"

        for run in title_slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(46)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
        for run in subtitle.text_frame.paragraphs[0].runs:
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(230, 230, 230)

        for idx in range(1, slide_count):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide_title = slide.shapes.title
            heading = f"{topic} • Part {idx}"
            if idx == 1:
                heading = f"What is {topic}?"
            elif idx == 2:
                heading = f"Key Applications of {topic}"
            elif idx == 3:
                heading = f"Benefits of {topic}"
            elif idx == slide_count - 1:
                heading = f"Future of {topic}"
            slide_title.text = heading
            for run in slide_title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(34)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

            body = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.2), Inches(4.8))
            body_frame = body.text_frame
            body_frame.word_wrap = True
            body_frame.clear()

            for bullet_idx in range(3):
                point_index = (idx * 3 + bullet_idx) % len(points)
                bullet_text = str(points[point_index]).strip()[:170]
                if bullet_idx == 0:
                    paragraph = body_frame.paragraphs[0]
                else:
                    paragraph = body_frame.add_paragraph()
                paragraph.text = bullet_text
                paragraph.level = 0
                paragraph.space_after = Pt(8)
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(240, 240, 240)

            if images:
                image_path = images[(idx - 1) % len(images)]
                try:
                    slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(1.7), width=Inches(5.8), height=Inches(3.8))
                except Exception:
                    pass

        output_name = f"{self._slugify_name(topic) or 'presentation'}_{time.strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = folder / output_name
        prs.save(str(output_path))
        self._open_external_path(str(output_path))
        return (
            f"Created a {slide_count}-slide PowerPoint on '{topic}' at {output_path}. "
            f"Downloaded {len(images)} topic images for design."
        )

    def _handle_image_download_request(self, text: str) -> str | None:
        raw = str(text or "").strip()
        lower = raw.lower()
        wants_image_download = (
            any(token in lower for token in ("image", "images", "photo", "photos", "wallpaper"))
            and any(token in lower for token in ("download", "save", "get"))
        )
        if not wants_image_download:
            return None
        topic = self._extract_topic(raw)
        if topic.lower() == "ai":
            infer = re.sub(
                r"\b(download|save|get|me|a|an|the|good|best|high quality|4k|hd|wallpaper|image|images|photo|photos|in|to|folder|downloads?|desktop|documents?)\b",
                " ",
                raw,
                flags=re.IGNORECASE,
            )
            infer = re.sub(r"\s+", " ", infer).strip(" .,:;\"'")
            if infer:
                topic = infer
        count = 6
        hit = re.search(r"\b(\d{1,2})\s+(?:images?|photos?|wallpapers?)\b", lower)
        if hit:
            try:
                count = max(1, min(20, int(hit.group(1))))
            except Exception:
                pass
        folder = self._resolve_export_folder(raw, default_name=f"{topic} Images")
        files = self._download_topic_images(topic, folder, count=count)
        self._open_external_path(str(folder))
        if not files:
            return "I couldn't download images right now. Please retry in a moment."
        return f"Downloaded {len(files)} images for '{topic}' in {folder}."

    def _extract_windows_path(self, text: str) -> str:
        raw = str(text or "")
        quoted = re.search(r"[\"']([a-zA-Z]:\\[^\"']+)[\"']", raw)
        if quoted:
            return quoted.group(1).strip()
        loose = re.search(r"([a-zA-Z]:\\[^\r\n]+?\.(?:png|jpg|jpeg|bmp|webp))", raw, flags=re.IGNORECASE)
        if loose:
            return loose.group(1).strip()
        return ""

    def _set_wallpaper(self, path_obj: Path) -> bool:
        if sys.platform != "win32" or ctypes is None:
            return False
        try:
            return bool(ctypes.windll.user32.SystemParametersInfoW(20, 0, str(path_obj), 3))
        except Exception:
            return False

    def _latest_image_candidate(self) -> Path | None:
        roots = [Path.home() / "Pictures", Path.home() / "Downloads", Path.home() / "Desktop"]
        latest = None
        latest_time = 0.0
        for root in roots:
            if not root.exists():
                continue
            try:
                for entry in root.glob("*"):
                    if not entry.is_file():
                        continue
                    if entry.suffix.lower() not in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
                        continue
                    mtime = entry.stat().st_mtime
                    if mtime > latest_time:
                        latest = entry
                        latest_time = mtime
            except Exception:
                continue
        return latest

    def _handle_wallpaper_command(self, text: str) -> str | None:
        raw = str(text or "").strip()
        lower = raw.lower()
        if "wallpaper" not in lower and "background" not in lower:
            return None
        if not any(token in lower for token in ("set", "change", "make", "apply")):
            return None
        if sys.platform != "win32":
            return "Wallpaper control is currently supported on Windows desktop only."

        candidate_path = Path(self._extract_windows_path(raw)).expanduser() if self._extract_windows_path(raw) else None
        if (not candidate_path or not candidate_path.exists()) and ("download" in lower or "google" in lower or "image" in lower):
            topic = self._extract_topic(raw)
            folder = self._resolve_export_folder(raw, default_name=f"{topic} Wallpapers")
            files = self._download_topic_images(topic, folder, count=1)
            if files:
                candidate_path = files[0]
        if (not candidate_path or not candidate_path.exists()) and ("file explorer" in lower or "explorer" in lower):
            candidate_path = self._latest_image_candidate()

        if not candidate_path or not candidate_path.exists():
            return "I couldn't find an image path. Provide a full image path, or ask me to download a wallpaper by topic."
        ok = self._set_wallpaper(candidate_path)
        if not ok:
            return "I found the image but couldn't set it as wallpaper."
        return f"Wallpaper updated from: {candidate_path}"

    def _set_windows_color_mode(self, dark_mode: bool) -> bool:
        if sys.platform != "win32" or winreg is None:
            return False
        try:
            key_path = r"Software\Microsoft\Windows\CurrentVersion\Themes\Personalize"
            with winreg.CreateKey(winreg.HKEY_CURRENT_USER, key_path) as key:
                value = 0 if dark_mode else 1
                winreg.SetValueEx(key, "AppsUseLightTheme", 0, winreg.REG_DWORD, value)
                winreg.SetValueEx(key, "SystemUsesLightTheme", 0, winreg.REG_DWORD, value)
            return True
        except Exception:
            return False

    def _handle_system_settings_command(self, text: str) -> str | None:
        raw = str(text or "").strip()
        lower = raw.lower()
        if not any(token in lower for token in ("setting", "settings", "wifi", "bluetooth", "display", "sound", "volume", "dark mode", "light mode", "wallpaper")):
            return None

        if "dark mode" in lower:
            if self._set_windows_color_mode(True):
                return "Switched system theme to dark mode."
            return "I couldn't switch dark mode automatically."
        if "light mode" in lower:
            if self._set_windows_color_mode(False):
                return "Switched system theme to light mode."
            return "I couldn't switch light mode automatically."

        if sys.platform != "win32":
            return None

        uri = ""
        if "bluetooth" in lower:
            uri = "ms-settings:bluetooth"
        elif "wifi" in lower or "network" in lower:
            uri = "ms-settings:network-wifi"
        elif "display" in lower or "resolution" in lower:
            uri = "ms-settings:display"
        elif "sound" in lower or "volume" in lower:
            uri = "ms-settings:sound"
        elif "wallpaper" in lower or "background" in lower:
            uri = "ms-settings:personalization-background"
        elif "update" in lower:
            uri = "ms-settings:windowsupdate"
        elif "privacy" in lower:
            uri = "ms-settings:privacy"
        elif "apps" in lower:
            uri = "ms-settings:appsfeatures"
        elif "setting" in lower:
            uri = "ms-settings:"
        if uri:
            opened = self._open_external_url(uri)
            if opened:
                return f"Opened system settings: {uri}"
        return None

    def _handle_optimize_command(self, text: str) -> str | None:
        raw = str(text or "").strip()
        lower = raw.lower()
        optimize_intent = any(
            token in lower
            for token in (
                "optimize",
                "optimise",
                "optimize pc",
                "optimise pc",
                "clean memory",
                "clear memory",
                "free ram",
                "memory cleanup",
                "memory clean",
            )
        )
        if not optimize_intent:
            return None

        root = Path(__file__).resolve().parent
        memreduct_exe = root / "optimization" / "bin" / "memreduct.exe"
        if not memreduct_exe.exists():
            memreduct_exe = root / "optimization" / "memreduct.exe"
        if not memreduct_exe.exists():
            memreduct_exe = root / "optimization" / "bin" / "memreduct-install" / "memreduct.exe"
        if not memreduct_exe.exists():
            return (
                "Mem Reduct binary not found. Place memreduct.exe in optimization\\bin "
                "or install Mem Reduct into optimization\\bin\\memreduct-install (admin required)."
            )

        try:
            subprocess.Popen(
                [str(memreduct_exe), "-clean"],
                cwd=str(memreduct_exe.parent),
                creationflags=subprocess.CREATE_NO_WINDOW,
            )
            return "Optimization started: memory cleanup is running now."
        except Exception as exc:
            return f"Optimization failed: {exc}"

    def _slugify_name(self, value: str) -> str:
        cleaned = re.sub(r"[^a-zA-Z0-9_-]+", "_", str(value or "").strip().lower()).strip("_")
        return cleaned[:64] if cleaned else ""

    def _infer_project_language(self, text: str) -> str:
        lower = str(text or "").lower()
        if "html" in lower and not any(k in lower for k in ("react", "next", "vue", "angular", "svelte")):
            return "html"
        if any(k in lower for k in ("python", "flask", "django", "fastapi")):
            return "python"
        if any(k in lower for k in ("react", "next", "node", "express", "vite", "tailwind", "javascript", "typescript", "web")):
            return "javascript"
        return "javascript"

    def _write_text_file(self, path_obj: Path, content: str):
        path_obj.parent.mkdir(parents=True, exist_ok=True)
        path_obj.write_text(content, encoding="utf-8")

    def _create_auth_template(self, target_dir: Path, brand: str):
        title = brand or "Brahma"
        html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{title} Auth</title>
  <link rel="stylesheet" href="./styles.css" />
</head>
<body>
  <main class="auth-shell">
    <section class="card">
      <p class="kicker">{title.upper()} ACCESS</p>
      <h1>Welcome back</h1>
      <p class="sub">Sign in to continue to your workspace.</p>
      <form id="authForm" class="form">
        <label>Email</label>
        <input type="email" id="email" placeholder="you@example.com" required />
        <label>Password</label>
        <input type="password" id="password" placeholder="Enter password" required />
        <button type="submit">Sign In</button>
      </form>
      <div id="msg" class="msg"></div>
    </section>
  </main>
  <script src="./app.js"></script>
</body>
</html>
"""
        css = """* { box-sizing: border-box; }
body {
  margin: 0;
  min-height: 100vh;
  font-family: 'Segoe UI', system-ui, sans-serif;
  background: radial-gradient(circle at 30% 10%, #14365c 0%, #081425 45%, #050d18 100%);
  color: #dff7ff;
}
.auth-shell {
  min-height: 100vh;
  display: grid;
  place-items: center;
  padding: 24px;
}
.card {
  width: min(420px, 100%);
  background: rgba(9, 20, 36, 0.88);
  border: 1px solid rgba(66, 232, 255, 0.28);
  border-radius: 20px;
  padding: 28px;
  box-shadow: 0 12px 50px rgba(0, 0, 0, 0.45);
}
.kicker {
  margin: 0 0 8px;
  color: #5de9ff;
  letter-spacing: 0.18em;
  font-size: 0.72rem;
}
h1 { margin: 0; font-size: 2rem; }
.sub { margin: 10px 0 22px; color: #9dbad1; }
.form { display: grid; gap: 10px; }
label { color: #9ad7ea; font-size: 0.92rem; }
input {
  width: 100%;
  padding: 12px 13px;
  border-radius: 12px;
  border: 1px solid rgba(96, 139, 165, 0.35);
  background: rgba(5, 18, 33, 0.86);
  color: #e6fbff;
}
input:focus { outline: 2px solid rgba(66, 232, 255, 0.6); border-color: transparent; }
button {
  margin-top: 6px;
  border: 0;
  border-radius: 12px;
  padding: 12px 14px;
  font-size: 1rem;
  font-weight: 600;
  color: #041625;
  background: linear-gradient(90deg, #32dfff 0%, #6eb8ff 100%);
  cursor: pointer;
}
button:hover { filter: brightness(1.07); }
.msg { margin-top: 12px; color: #8fe8a0; min-height: 1.2em; }
"""
        js = """const form = document.getElementById('authForm');
const msg = document.getElementById('msg');
form.addEventListener('submit', (e) => {
  e.preventDefault();
  msg.textContent = 'Signed in (demo). Connect backend auth API to make this live.';
});
"""
        self._write_text_file(target_dir / "index.html", html)
        self._write_text_file(target_dir / "styles.css", css)
        self._write_text_file(target_dir / "app.js", js)
        self._write_text_file(target_dir / "README.md", f"# {title} Auth Page\n\nGenerated in fallback mode.")

    def _create_basic_web_template(self, target_dir: Path, brand: str):
        title = brand or "Brahma Web"
        html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{title}</title>
  <link rel="stylesheet" href="./styles.css" />
</head>
<body>
  <header>
    <h1>{title}</h1>
    <p>Modern starter website generated by Brahma.</p>
  </header>
  <main>
    <section class="card">
      <h2>Fast Start</h2>
      <p>Edit <code>index.html</code>, <code>styles.css</code>, and <code>app.js</code>.</p>
    </section>
  </main>
  <script src="./app.js"></script>
</body>
</html>
"""
        css = """body{margin:0;font-family:Segoe UI,system-ui,sans-serif;background:#081426;color:#e7f7ff}
header{padding:56px 24px;background:linear-gradient(120deg,#0b2038,#0e3557)}
h1{margin:0 0 8px}
main{padding:24px}
.card{background:#0d1f35;border:1px solid #1f4369;border-radius:14px;padding:18px;max-width:680px}
"""
        js = "console.log('Brahma starter website ready.');\n"
        self._write_text_file(target_dir / "index.html", html)
        self._write_text_file(target_dir / "styles.css", css)
        self._write_text_file(target_dir / "app.js", js)
        self._write_text_file(target_dir / "README.md", f"# {title}\n\nGenerated in fallback mode.")

    def _create_notes_template(self, target_dir: Path, brand: str):
        title = brand or "Brahma Notes"
        html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{title}</title>
  <link rel="stylesheet" href="./styles.css" />
</head>
<body>
  <main class="wrap">
    <header>
      <h1>{title}</h1>
      <p>Local notes. Everything is saved in your browser on this device.</p>
    </header>
    <section class="composer">
      <input id="titleInput" placeholder="Note title" />
      <textarea id="bodyInput" rows="6" placeholder="Write your note..."></textarea>
      <button id="saveBtn">Save Note</button>
    </section>
    <section>
      <h2>Saved Notes</h2>
      <div id="list" class="list"></div>
    </section>
  </main>
  <script src="./app.js"></script>
</body>
</html>
"""
        css = """*{box-sizing:border-box}body{margin:0;font-family:Segoe UI,system-ui,sans-serif;background:#071325;color:#e9f8ff}
.wrap{max-width:900px;margin:0 auto;padding:26px}
header p{color:#98b8cf}
.composer{display:grid;gap:10px;background:#0b2038;border:1px solid #1c4268;border-radius:14px;padding:14px}
input,textarea{width:100%;padding:10px 12px;border:1px solid #265072;border-radius:10px;background:#0a1b30;color:#e9f8ff}
button{padding:10px 14px;border:0;border-radius:10px;background:#45dcff;color:#05223a;font-weight:700;cursor:pointer}
.list{display:grid;gap:10px}
.note{background:#0a1d34;border:1px solid #214666;border-radius:12px;padding:12px}
.note h3{margin:0 0 6px}
.muted{color:#89a8bf;font-size:.85rem}
"""
        js = """const KEY='brahma_local_notes_v1';
const titleInput=document.getElementById('titleInput');
const bodyInput=document.getElementById('bodyInput');
const saveBtn=document.getElementById('saveBtn');
const list=document.getElementById('list');

function load(){try{return JSON.parse(localStorage.getItem(KEY)||'[]')}catch{return[]}}
function store(notes){localStorage.setItem(KEY,JSON.stringify(notes))}
function render(){
  const notes=load();
  list.innerHTML='';
  if(!notes.length){list.innerHTML='<p class=\"muted\">No notes yet.</p>';return}
  notes.forEach((n,idx)=>{
    const el=document.createElement('article');
    el.className='note';
    el.innerHTML=`<h3>${n.title||'Untitled'}</h3><p>${n.body||''}</p><div class=\"muted\">Saved: ${new Date(n.ts).toLocaleString()}</div><button data-i=\"${idx}\">Delete</button>`;
    list.appendChild(el);
  });
}
saveBtn.addEventListener('click',()=>{
  const title=titleInput.value.trim();
  const body=bodyInput.value.trim();
  if(!title && !body) return;
  const notes=load();
  notes.unshift({title,body,ts:Date.now()});
  store(notes);
  titleInput.value='';bodyInput.value='';
  render();
});
list.addEventListener('click',(e)=>{
  const btn=e.target.closest('button[data-i]');
  if(!btn) return;
  const i=Number(btn.dataset.i);
  const notes=load();
  notes.splice(i,1);
  store(notes);
  render();
});
render();
"""
        self._write_text_file(target_dir / "index.html", html)
        self._write_text_file(target_dir / "styles.css", css)
        self._write_text_file(target_dir / "app.js", js)
        self._write_text_file(target_dir / "README.md", f"# {title}\n\nLocalStorage notes website fallback template.")

    def _create_article_template(self, target_dir: Path, brand: str, topic: str):
        title = brand or "Brahma Articles"
        safe_topic = topic or "AI"
        html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{title} - {safe_topic}</title>
  <link rel="stylesheet" href="./styles.css" />
</head>
<body>
  <main class="layout">
    <article class="article">
      <p class="kicker">{title.upper()}</p>
      <h1>{safe_topic}: A Practical Overview</h1>
      <p class="lead">This article page was generated in fallback mode and is ready to edit.</p>
      <h2>What It Is</h2>
      <p>{safe_topic} impacts products, learning, and daily workflows. Teams now prioritize usable systems over demos.</p>
      <h2>Why It Matters</h2>
      <p>It improves speed and consistency, but outcomes still depend on quality data, clear goals, and human review.</p>
      <h2>What Comes Next</h2>
      <p>Future progress will combine better models, stronger tooling, and practical governance for reliability and safety.</p>
    </article>
  </main>
</body>
</html>
"""
        css = """body{margin:0;font-family:Georgia,'Times New Roman',serif;background:#f5f3ee;color:#222}
.layout{max-width:920px;margin:0 auto;padding:42px 20px}
.article{background:#fff;border:1px solid #ded8cb;border-radius:14px;padding:28px 26px;box-shadow:0 10px 30px rgba(0,0,0,.06)}
.kicker{margin:0;color:#33506c;letter-spacing:.16em;font-size:.76rem}
h1{font-size:2.3rem;line-height:1.1;margin:.4rem 0 1rem}
.lead{font-size:1.1rem;color:#4b4b4b}
h2{margin-top:1.5rem}
p{line-height:1.75}
"""
        self._write_text_file(target_dir / "index.html", html)
        self._write_text_file(target_dir / "styles.css", css)
        self._write_text_file(target_dir / "README.md", f"# {title}\n\nArticle website fallback template on topic: {safe_topic}")

    def _fallback_project_build(self, text: str, target_dir: Path, reason: str = "") -> str:
        lower = str(text or "").lower()
        brand = self._extract_brand_name(text) or target_dir.name
        topic = self._extract_topic(text)
        reason_note = self._friendly_builder_reason(reason)
        if any(token in lower for token in ("note", "notes", "notepad")):
            self._create_notes_template(target_dir, brand)
            self._open_external_path(str(target_dir / "index.html"))
            return (
                f"Created notes website in {target_dir} (fallback template). "
                f"{reason_note}"
            )
        if any(token in lower for token in ("article", "blog", "news", "history")):
            self._create_article_template(target_dir, brand, topic)
            self._open_external_path(str(target_dir / "index.html"))
            return (
                f"Created article website in {target_dir} (fallback template). "
                f"{reason_note}"
            )
        if any(token in lower for token in ("auth", "login", "sign in", "signin", "signup", "register")):
            self._create_auth_template(target_dir, brand)
            self._open_external_path(str(target_dir / "index.html"))
            return (
                f"Created auth website in {target_dir} (fallback template). "
                f"{reason_note}"
            )
        self._create_basic_web_template(target_dir, brand)
        self._open_external_path(str(target_dir / "index.html"))
        return (
            f"Created starter website in {target_dir} (fallback template). "
            f"{reason_note}"
        )

    def _normalize_routine_key(self, value: str) -> str:
        key = str(value or "").strip().lower()
        key = key.replace("rountine", "routine")
        key = re.sub(r"[^a-z0-9]+", "_", key).strip("_")
        aliases = {
            "morning": "morning_setup",
            "morning_routine": "morning_setup",
            "morning_setup_routine": "morning_setup",
            "work": "work_mode",
            "work_mode_routine": "work_mode",
            "work_routine": "work_mode",
            "gaming": "gaming_mode",
            "game_mode": "gaming_mode",
            "gaming_routine": "gaming_mode",
            "gaming_mode_routine": "gaming_mode",
        }
        return aliases.get(key, key)

    def _default_routines(self):
        return {
            "morning_setup": {
                "name": "Morning Setup",
                "steps": [
                    "open chrome",
                    "open file explorer",
                    "organize downloads",
                ],
            },
            "work_mode": {
                "name": "Work Mode",
                "steps": [
                    "open vscode",
                    "open chrome",
                    "turn on focus mode",
                ],
            },
            "gaming_mode": {
                "name": "Gaming Mode",
                "steps": [
                    "open steam",
                    "open discord",
                    "turn off notifications",
                ],
            },
        }

    def _execute_routine_step(self, step: str) -> str:
        line = str(step or "").strip()
        if not line:
            return "Skipped empty step."
        lower = line.lower()
        if re.search(r"\b(run|start|execute)\b.+\brou?tine\b", lower):
            return "Skipped nested routine command."
        direct = self.try_handle_direct_command(line)
        if direct:
            return direct
        result, _meta = self.execute_hybrid_command(line, force_openclaw=False)
        return str(result or "Done.")

    def _list_routines_text(self) -> str:
        listed = workflow_manager({"action": "list"}, player=self)
        if "no workflows saved yet" in listed.lower():
            return "No routines saved yet."
        return listed.replace("Saved workflows:", "Available routines:")

    def _run_routine(self, name_hint: str) -> str:
        key = self._normalize_routine_key(name_hint)
        defaults = self._default_routines()
        candidate_names = [key, key.replace("_", " ")]

        if key in defaults:
            display_name = defaults[key]["name"]
            show = workflow_manager({"action": "show", "name": display_name}, player=self)
            if "not found" in show.lower():
                workflow_manager(
                    {"action": "create", "name": display_name, "steps": defaults[key]["steps"]},
                    player=self,
                )
        else:
            listed = workflow_manager({"action": "list"}, player=self)
            names = []
            if ":" in listed:
                names = [item.strip() for item in listed.split(":", 1)[1].split(",") if item.strip()]
            if names:
                mapping = {self._normalize_routine_key(n): n for n in names}
                if key in mapping:
                    display_name = mapping[key]
                else:
                    close = difflib.get_close_matches(key, list(mapping.keys()), n=1, cutoff=0.65)
                    if close:
                        display_name = mapping[close[0]]
                    else:
                        display_name = ""
            else:
                display_name = ""
            if not display_name:
                return f"Routine '{name_hint}' not found. {self._list_routines_text()}"

        if key in defaults:
            display_name = defaults[key]["name"]
        run_result = workflow_manager(
            {"action": "run", "name": display_name, "step_runner": self._execute_routine_step},
            player=self,
        )
        self.refresh_sequences()
        return run_result

    def _handle_routine_command(self, text: str):
        raw = str(text or "").strip()
        lower = raw.lower().replace("rountine", "routine")
        if lower in {"list routines", "show routines", "list routine", "show routine"}:
            self.refresh_sequences()
            return self._list_routines_text()

        run_match = re.search(r"^(?:run|start|execute)\s+(.+?)\s*(?:routine)?\s*$", lower)
        if run_match and ("mode" in lower or "routine" in lower):
            target = run_match.group(1).strip()
            target = re.sub(r"\s+routine$", "", target).strip()
            return self._run_routine(target)
        if re.search(r"\brou?tine\b", lower) and any(token in lower for token in ("run ", "start ", "execute ")):
            target = re.sub(r"^(run|start|execute)\s+", "", lower).strip()
            target = re.sub(r"\s+routine$", "", target).strip()
            return self._run_routine(target)
        return None

    def _handle_project_build(self, text: str) -> str:
        workspace = self._current_project_workspace()
        if not workspace:
            return (
                "Project workspace is not selected yet. Open Project Workspace and choose a folder first, "
                "then ask me to create the website/app."
            )

        requested_name = self._extract_named_project(text)
        slug = self._slugify_name(requested_name)
        kind = "website" if "web" in text.lower() or "site" in text.lower() or "html" in text.lower() else "app"
        if not slug:
            slug = f"{kind}_{time.strftime('%Y%m%d_%H%M%S')}"
        target_dir = workspace / slug
        target_dir.mkdir(parents=True, exist_ok=True)

        from actions.dev_agent import dev_agent

        language = self._infer_project_language(text)
        try:
            build_result = dev_agent(
                parameters={
                    "description": text,
                    "language": language,
                    "project_root": str(target_dir),
                    "timeout": 45,
                },
                player=self,
                speak=self._speak_response,
            )
        except Exception as exc:
            return self._fallback_project_build(text, target_dir, reason=str(exc))

        result_text = str(build_result or "").strip()
        if not result_text:
            return self._fallback_project_build(text, target_dir, reason="Empty builder response")

        lowered = result_text.lower()
        if any(token in lowered for token in ("rate limit", "quota", "resourceexhausted", "planning failed", "could not")):
            return self._fallback_project_build(text, target_dir, reason=result_text)

        self._open_external_path(str(target_dir))
        return f"Created project in {target_dir}. {result_text}"

    def _resolve_download_dir(self, text: str) -> Path:
        lower = str(text or "").lower()
        home = Path.home()
        if "desktop" in lower:
            target = home / "Desktop"
        elif "document" in lower:
            target = home / "Documents"
        else:
            target = home / "Downloads"
        target.mkdir(parents=True, exist_ok=True)
        return target

    def _is_youtube_url(self, url: str) -> bool:
        if not url:
            return False
        lower = url.lower()
        return "youtube.com" in lower or "youtu.be" in lower

    def _is_spotify_url(self, url: str) -> bool:
        return bool(url and "spotify.com" in url.lower())

    def _open_youtube_non_popup(self, text: str) -> str:
        query = re.sub(r"\b(open|search|play|on|in)\b", " ", str(text or ""), flags=re.IGNORECASE)
        query = re.sub(r"\byoutube\b", " ", query, flags=re.IGNORECASE)
        query = re.sub(r"\s+", " ", query).strip()
        url = "https://www.youtube.com"
        if query:
            url = f"https://www.youtube.com/results?search_query={quote_plus(query)}"
        opened, _ = self._open_in_managed_browser(url, allow_external_fallback=True)
        if opened:
            return "Opened YouTube in your current browser window (no popup)."
        return "I couldn't open YouTube right now."

    def _open_spotify_non_popup(self, text: str) -> str:
        direct_url = self._extract_first_url(text)
        if self._is_spotify_url(direct_url):
            url = direct_url
        else:
            query = re.sub(r"\b(open|search|play|on|in)\b", " ", str(text or ""), flags=re.IGNORECASE)
            query = re.sub(r"\bspotify\b", " ", query, flags=re.IGNORECASE)
            query = re.sub(r"\s+", " ", query).strip()
            url = "https://open.spotify.com"
            if query:
                url = f"https://open.spotify.com/search/{quote_plus(query)}"
        opened, _ = self._open_in_managed_browser(url, allow_external_fallback=True)
        if opened:
            return "Opened Spotify in your current browser window (no popup)."
        return "I couldn't open Spotify right now."

    def _handle_browser_navigation(self, text: str):
        raw = str(text or "").strip()
        if not raw:
            return None
        lower = raw.lower()

        if any(token in lower for token in ("youtube", "spotify")):
            return None
        if any(token in lower for token in ("send email", "send mail", "reply email", "delete email", "compose email")):
            return None

        navigation_terms = ("open ", "go to", "goto", "navigate", "visit", "browse", "search")
        if not any(term in lower for term in navigation_terms):
            return None

        direct_url = self._extract_first_url(raw)
        if direct_url:
            opened, _ = self._open_in_managed_browser(direct_url, allow_external_fallback=True)
            return "Opened that link in your current browser window." if opened else "I couldn't open that link right now."

        if "gmail" in lower and not any(token in lower for token in ("send", "reply", "delete", "compose")):
            opened, _ = self._open_in_managed_browser("https://mail.google.com", allow_external_fallback=True)
            return "Opened Gmail in your current browser window." if opened else "I couldn't open Gmail right now."
        if "discord" in lower:
            opened, _ = self._open_in_managed_browser("https://discord.com/app", allow_external_fallback=True)
            return "Opened Discord in your current browser window." if opened else "I couldn't open Discord right now."

        if any(token in lower for token in ("open chrome", "open browser", "open google", "go to google", "search google")):
            query = ""
            if "search" in lower:
                query = re.sub(
                    r"\b(open|go|to|goto|navigate|visit|browse|search|for|on|in|chrome|browser|google)\b",
                    " ",
                    raw,
                    flags=re.IGNORECASE,
                )
                query = re.sub(r"\s+", " ", query).strip()
            target = "https://www.google.com"
            if query:
                target = f"https://www.google.com/search?q={quote_plus(query)}"
            opened, _ = self._open_in_managed_browser(target, allow_external_fallback=True)
            return "Opened in your current browser window." if opened else "I couldn't open that in browser right now."

        return None

    def _handle_browser_agent_command(self, text: str) -> str | None:
        raw = str(text or "").strip()
        if not raw:
            return None
        lower = raw.lower()
        triggers = ("browse", "surf", "agentic", "browser agent", "web agent", "multi-step browser")
        if not any(token in lower for token in triggers):
            return None
        try:
            result = self.browser_agent_bridge.run_command(raw)
        except Exception as exc:
            return f"Browser agent failed: {exc}"
        if not result.get("ok"):
            return f"Browser agent error: {result.get('error') or 'unknown error'}"
        steps = result.get("results") or []
        summary = steps[-1] if steps else "Browser agent completed."
        return f"{summary}"

    def _lookup_steam_app(self, query: str):
        term = str(query or "").strip()
        if not term:
            return None
        try:
            api_url = f"https://store.steampowered.com/api/storesearch/?term={quote_plus(term)}&l=english&cc=us"
            req = Request(api_url, headers={"User-Agent": "BrahmaAI/4.0"})
            with urlopen(req, timeout=12) as response:
                payload = json.loads(response.read().decode("utf-8", errors="replace"))
            items = payload.get("items") or []
            if not items:
                return None
            normalized_query = re.sub(r"[^a-z0-9]+", " ", term.lower()).strip()

            def _score(item):
                name = str((item or {}).get("name") or "")
                normalized_name = re.sub(r"[^a-z0-9]+", " ", name.lower()).strip()
                if not normalized_name:
                    return -1.0
                ratio = difflib.SequenceMatcher(None, normalized_query, normalized_name).ratio()
                if normalized_query == normalized_name:
                    ratio += 1.0
                elif normalized_query in normalized_name:
                    ratio += 0.45
                elif normalized_name in normalized_query:
                    ratio += 0.25
                return ratio

            best = max(items, key=_score)
            appid = (best or {}).get("id")
            name = (best or {}).get("name") or term
            if appid is None:
                return None
            return {"appid": int(appid), "name": str(name)}
        except Exception:
            return None

    def _extract_steam_appid(self, text: str):
        raw = str(text or "")
        patterns = [
            r"store\.steampowered\.com/app/(\d+)",
            r"\b(?:app\s*id|appid|game\s*id|id)\s*[-:=]?\s*(\d{3,10})\b",
        ]
        for pattern in patterns:
            match = re.search(pattern, raw, flags=re.IGNORECASE)
            if match:
                try:
                    return int(match.group(1))
                except Exception:
                    return None
        return None

    def _clean_steam_query(self, text: str) -> str:
        raw = str(text or "")
        raw = re.sub(r"https?://\S+", " ", raw, flags=re.IGNORECASE)
        lowered = raw.lower().replace("ets2", "ets 2")
        if "ets 2" in lowered:
            lowered = lowered.replace("ets 2", "euro truck simulator 2")
        lowered = re.sub(r"\b[a-z]\s*drive\b", " ", lowered, flags=re.IGNORECASE)
        lowered = re.sub(r"\bdrive\s*[a-z]\b", " ", lowered, flags=re.IGNORECASE)
        lowered = re.sub(r"\b[a-z]:\\?\b", " ", lowered, flags=re.IGNORECASE)
        cleaned = re.sub(
            r"\b(download|install|get|from|on|in|using|steam|game|games|library|app|please|me|its|it's|id|game id|app id|folder|disk|location|path|into|to)\b",
            " ",
            lowered,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(r"\b\d{3,10}\b", " ", cleaned)
        cleaned = re.sub(r"[^a-z0-9 ]+", " ", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s+", " ", cleaned).strip(" .,:;")
        return cleaned

    def _steam_known_appid(self, query: str):
        normalized = re.sub(r"\s+", " ", str(query or "").strip().lower())
        if not normalized:
            return None
        known = {
            "euro truck simulator 2": 227300,
            "ets 2": 227300,
            "counter strike 2": 730,
            "cs2": 730,
            "dota 2": 570,
            "gta v": 271590,
            "grand theft auto v": 271590,
            "pubg": 578080,
            "elden ring": 1245620,
            "forza horizon 5": 1551360,
        }
        if normalized in known:
            return known[normalized]
        for name, appid in known.items():
            if normalized in name or name in normalized:
                return appid
        return None

    def _handle_steam_download(self, text: str) -> str:
        raw = str(text or "")
        explicit_appid = self._extract_steam_appid(raw)
        if explicit_appid:
            opened = self._open_external_url(f"steam://install/{explicit_appid}")
            if opened:
                return f"Sent install command to Steam for AppID {explicit_appid}."
            self._open_external_url(f"https://store.steampowered.com/app/{explicit_appid}/")
            return f"Found AppID {explicit_appid}. Opened store page because Steam install URI failed."

        query = self._clean_steam_query(raw)
        if not query:
            self._open_external_url("steam://open/main")
            return "Steam opened. Tell me the exact game name to start install."

        known_appid = self._steam_known_appid(query)
        if known_appid:
            opened = self._open_external_url(f"steam://install/{known_appid}")
            if opened:
                return f"Sent install command to Steam for '{query}' (AppID: {known_appid})."

        found = self._lookup_steam_app(query)
        if found:
            appid = found["appid"]
            name = found["name"]
            opened = self._open_external_url(f"steam://install/{appid}")
            if opened:
                return f"Sent install command to Steam for '{name}' (AppID: {appid})."
            self._open_external_url(f"https://store.steampowered.com/app/{appid}/")
            return f"Found '{name}' (AppID: {appid}). Opened store page because Steam install URI failed."

        self._open_external_url(f"https://store.steampowered.com/search/?term={quote_plus(query)}")
        return f"I couldn't resolve an exact Steam app for '{query}'. Opened Steam search results."

    def _handle_youtube_download(self, text: str) -> str:
        url = self._extract_first_url(text)
        source = ""
        if self._is_youtube_url(url):
            source = url
        else:
            query = re.sub(
                r"\b(download|save|youtube|video|music|song|audio|mp3|as|to|in|on|from)\b",
                " ",
                str(text or ""),
                flags=re.IGNORECASE,
            )
            query = re.sub(r"\s+", " ", query).strip()
            if not query:
                self._open_youtube_non_popup(text)
                return (
                    "Opened YouTube without popup. For direct download, provide a YouTube URL or a search phrase "
                    "for content you own or have permission to save."
                )
            source = f"ytsearch1:{query}"

        check = subprocess.run(
            [sys.executable, "-m", "yt_dlp", "--version"],
            capture_output=True,
            text=True,
            timeout=20,
        )
        if check.returncode != 0:
            install = subprocess.run(
                [sys.executable, "-m", "pip", "install", "--disable-pip-version-check", "-q", "yt-dlp"],
                capture_output=True,
                text=True,
                timeout=180,
            )
            if install.returncode != 0:
                err = (install.stderr or install.stdout or "").strip()
                if len(err) > 220:
                    err = err[:220] + "..."
                return (
                    "I couldn't install `yt-dlp` automatically right now. "
                    f"Install it once with: python -m pip install yt-dlp. Details: {err or 'install failed'}"
                )

        target_dir = self._resolve_download_dir(text)
        wants_audio = any(token in text.lower() for token in ("music", "song", "audio", "mp3"))
        outtmpl = str(target_dir / "%(title)s.%(ext)s")
        cmd = [sys.executable, "-m", "yt_dlp", "--no-playlist", "-o", outtmpl]
        if wants_audio:
            cmd.extend(["-f", "bestaudio/best"])
        else:
            cmd.extend(["-f", "bv*+ba/b"])
        cmd.append(source)
        creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0) if sys.platform == "win32" else 0
        subprocess.Popen(
            cmd,
            cwd=str(target_dir),
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            creationflags=creation_flags,
        )
        media_type = "audio" if wants_audio else "video"
        return (
            f"Started YouTube {media_type} download in background (no popup). "
            f"Target folder: {target_dir}"
        )

    def _extract_email_address(self, text: str) -> str:
        hit = re.search(r"[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}", str(text or ""), flags=re.IGNORECASE)
        return hit.group(0).strip() if hit else ""

    def _gmail_search_latest(self, mailbox, from_hint: str = "", subject_hint: str = ""):
        criteria = ["ALL"]
        if from_hint:
            criteria.extend(["FROM", f'"{from_hint}"'])
        if subject_hint:
            criteria.extend(["SUBJECT", f'"{subject_hint}"'])
        status, data = mailbox.search(None, *criteria)
        if status != "OK":
            return None
        ids = (data[0] or b"").split()
        if not ids:
            return None
        return ids[-1]

    def _handle_gmail_command(self, text: str):
        raw = str(text or "").strip()
        lower = raw.lower()
        if not any(token in lower for token in ("gmail", "email", "mail")):
            return None
        intent = any(
            token in lower
            for token in ("open", "start", "go to", "goto", "compose", "send", "reply", "delete")
        )
        if not intent:
            return None

        opened, _ = self._open_in_managed_browser("https://mail.google.com", allow_external_fallback=True)
        if not opened:
            return "I couldn't open Gmail in your current browser window."

        # Browser-control mode (keyboard/mouse style) instead of SMTP/IMAP credentials.
        if any(token in lower for token in ("compose", "send")):
            return "Opened Gmail in your current browser window and switched to browser-control mode for compose/send actions."
        if "reply" in lower:
            return "Opened Gmail in your current browser window and switched to browser-control mode for reply."
        if "delete" in lower:
            return "Opened Gmail in your current browser window and switched to browser-control mode for delete."
        return "Opened Gmail in your current browser window."

    def _handle_discord_command(self, text: str):
        raw = str(text or "").strip()
        lower = raw.lower()
        if "discord" not in lower:
            return None
        intent = any(
            token in lower
            for token in ("open", "start", "go to", "goto", "send", "message", "dm", "reply", "create")
        )
        if intent:
            opened, _ = self._open_in_managed_browser("https://discord.com/app", allow_external_fallback=True)
            if opened:
                return "Opened Discord in your current browser window and switched to browser-control mode."
            return "I couldn't open Discord in your current browser window."

        if "list" in lower and "server" in lower:
            if not self._discord_token():
                return "Discord bot token is not configured. Use 'open discord in current browser' to control it with browser automation."
            guilds = self._discord_list_guilds()
            if not guilds:
                return "No Discord servers are available for this bot."
            names = ", ".join(str(g.get("name") or g.get("id")) for g in guilds[:30])
            return f"Discord servers: {names}"

        if "create" in lower and "server" in lower:
            if not self._discord_token():
                return "Discord bot token is not configured. Use browser-control mode for Discord."
            server_name = self._discord_extract_server_name(raw) or "Brahma Server"
            payload = self._discord_api_request("POST", "/guilds", payload={"name": server_name})
            created_name = str(payload.get("name") or server_name)
            created_id = str(payload.get("id") or "")
            return f"Created Discord server '{created_name}'{f' (id: {created_id})' if created_id else ''}."

        if "create" in lower and "channel" in lower:
            if not self._discord_token():
                return "Discord bot token is not configured. Use browser-control mode for Discord."
            server_name = self._discord_extract_server_name(raw)
            channel_name = self._discord_extract_channel_name(raw)
            if not channel_name:
                return "Please provide channel name, for example: create discord channel #announcements in server My Server."
            guild = self._discord_resolve_guild(server_name)
            if not guild:
                return "I could not find that Discord server for channel creation."
            safe_channel_name = re.sub(r"[^a-z0-9_-]+", "-", channel_name.lower()).strip("-") or "new-channel"
            created = self._discord_api_request(
                "POST",
                f"/guilds/{guild.get('id')}/channels",
                payload={"name": safe_channel_name, "type": 0},
            )
            return f"Created channel #{created.get('name') or safe_channel_name} in {guild.get('name')}."

        if any(token in lower for token in ("offline", "auto reply", "auto-reply")):
            if not self._discord_token():
                return "Discord bot token is not configured. Use browser-control mode for Discord."
            disable = any(token in lower for token in ("disable", "turn off", "stop"))
            if disable:
                self.discord_settings["auto_reply_enabled"] = False
                self._save_discord_settings()
                return "Discord offline auto-reply has been disabled."
            server_name = self._discord_extract_server_name(raw)
            channel_ids = []
            if server_name:
                guild = self._discord_resolve_guild(server_name)
                if guild:
                    channels = self._discord_list_text_channels(str(guild.get("id") or ""))
                    channel_ids = [str(ch.get("id")) for ch in channels[:20] if str(ch.get("id") or "").strip()]
            reply_text = self._discord_extract_message_body(raw) or self.discord_settings.get("auto_reply_message") or (
                "I am offline right now. I will get back to you soon."
            )
            self.discord_settings["auto_reply_message"] = reply_text
            self.discord_settings["auto_reply_enabled"] = True
            if channel_ids:
                self.discord_settings["auto_reply_channel_ids"] = channel_ids
            self._save_discord_settings()
            if channel_ids:
                return (
                    f"Discord offline auto-reply enabled for {len(channel_ids)} channels "
                    f"in '{server_name}' with message: {reply_text}"
                )
            return f"Discord offline auto-reply enabled. Message: {reply_text}"

        send_intent = any(token in lower for token in ("send", "message", "dm", "direct message"))
        if send_intent and ("dm" in lower or "direct message" in lower):
            if not self._discord_token():
                return "Discord bot token is not configured. Use browser-control mode for Discord."
            user_match = re.search(r"<@!?(\d{15,24})>", raw) or re.search(r"\b(?:user|id)\s*[:=-]?\s*(\d{15,24})\b", raw)
            user_id = user_match.group(1).strip() if user_match else ""
            if not user_id:
                return "For Discord DM, provide user id, e.g. send discord dm to user id 123456789012345678 message hi."
            body = self._discord_extract_message_body(raw)
            if not body:
                return "Please include the DM message text."
            self._discord_send_dm(user_id, body)
            return f"Sent Discord DM to user id {user_id}."

        if send_intent:
            if not self._discord_token():
                return "Discord bot token is not configured. Use browser-control mode for Discord."
            server_name = self._discord_extract_server_name(raw)
            if not server_name:
                return "Please include a Discord server name, for example: send discord message in server My Server channel general ..."
            guild = self._discord_resolve_guild(server_name)
            if not guild:
                return f"I couldn't find Discord server '{server_name}'."
            channel_hint = self._discord_extract_channel_name(raw)
            channel = self._discord_resolve_channel(str(guild.get("id") or ""), channel_hint)
            if not channel:
                return f"I couldn't find a text channel in '{guild.get('name')}'."
            body = self._discord_extract_message_body(raw)
            if not body:
                return "Please include the message text."
            self._discord_send_channel_message(str(channel.get("id")), body)
            return f"Sent message to #{channel.get('name')} in {guild.get('name')}."

        return (
            "Discord feature is ready. You can ask me to create a server, create a channel, send server/DM messages, "
            "or enable offline auto-reply."
        )

    def try_handle_direct_command(self, text: str):
        raw = str(text or "").strip()
        if not raw:
            return None
        lower = raw.lower()
        ppt_intent = (
            any(token in lower for token in ("ppt", "pptx", "powerpoint", "presentation", "slides"))
            and any(token in lower for token in ("create", "make", "generate", "build", "design"))
        )
        if ppt_intent:
            return self._handle_ppt_request(raw)

        spreadsheet_intent = (
            any(token in lower for token in ("spreadsheet", "spreadsheets", "excel", ".xlsx", "sheet"))
            and any(token in lower for token in ("create", "make", "generate", "build", "research"))
        )
        if spreadsheet_intent:
            return self._handle_spreadsheet_request(raw)

        wallpaper_result = self._handle_wallpaper_command(raw)
        if wallpaper_result:
            return wallpaper_result

        optimize_result = self._handle_optimize_command(raw)
        if optimize_result:
            return optimize_result

        image_download_result = self._handle_image_download_request(raw)
        if image_download_result:
            return image_download_result

        system_settings_result = self._handle_system_settings_command(raw)
        if system_settings_result:
            return system_settings_result

        try:
            gmail_result = self._handle_gmail_command(raw)
        except Exception as exc:
            return f"Gmail command failed: {exc}"
        if gmail_result:
            return gmail_result

        try:
            discord_result = self._handle_discord_command(raw)
        except Exception as exc:
            return f"Discord command failed: {exc}"
        if discord_result:
            return discord_result

        routine_result = self._handle_routine_command(raw)
        if routine_result:
            return routine_result
        browser_agent_result = self._handle_browser_agent_command(raw)
        if browser_agent_result:
            return browser_agent_result
        browser_navigation = self._handle_browser_navigation(raw)
        if browser_navigation:
            return browser_navigation
        build_intent = (
            any(token in lower for token in ("create", "build", "make", "develop", "generate", "scaffold"))
            and any(token in lower for token in ("website", "web app", "webapp", "site", "application", " app ", " app", "android app", "desktop app", "mobile app"))
        )
        if build_intent:
            return self._handle_project_build(raw)

        wants_download = any(token in lower for token in ("download", "save offline", "save to"))
        mentions_youtube = ("youtube" in lower) or ("youtu.be" in lower) or self._is_youtube_url(self._extract_first_url(raw))
        mentions_spotify = ("spotify" in lower) or self._is_spotify_url(self._extract_first_url(raw))
        wants_steam_install = ("steam" in lower) and any(token in lower for token in ("download", "install", "get"))

        if wants_steam_install:
            return self._handle_steam_download(raw)

        if mentions_youtube and wants_download:
            return self._handle_youtube_download(raw)

        if mentions_spotify and wants_download:
            self._open_spotify_non_popup(raw)
            return (
                "Opened Spotify without popup. Spotify media files are DRM-protected, so direct file download "
                "is only supported inside the official Spotify app using its offline feature."
            )

        if mentions_youtube:
            return self._open_youtube_non_popup(raw)
        if mentions_spotify:
            return self._open_spotify_non_popup(raw)
        return None

    def _run_brahma_multiaction(self, text: str):
        from agent.executor import AgentExecutor

        executor = AgentExecutor()
        return executor.execute(goal=text, speak=self._speak_response)

    def _speak_response(self, text: str):
        text = str(text or "").strip()
        if not text:
            return
        try:
            speak_text(text, settings=self._edge_only_voice_settings(self.voice_settings), player=self)
        except Exception:
            pass

    def _speak_response_async(self, text: str):
        text = str(text or "").strip()
        if not text:
            return
        threading.Thread(
            target=self._speak_response,
            args=(text,),
            daemon=True,
            name="BrahmaVoice",
        ).start()

    def execute_hybrid_command(self, text: str, force_openclaw: bool = False):
        effective_text = self._inject_project_context(text)
        route = route_command(text, advanced_mode=False)
        metadata = {
            "engine": route.engine,
            "confidence": route.confidence,
            "reason": route.reason,
        }

        if route.engine == "brahma_multiaction":
            return self._run_brahma_multiaction(effective_text), metadata

        return offline_assistant(effective_text, player=self), metadata

    def start_streamed_advanced_command(self, text: str):
        effective_text = self._inject_project_context(text)
        self.advanced_terminal_active = True
        self.status_text = "PROCESSING"
        self.write_log("[sys] Advanced mode routed to standard Brahma flow.")
        try:
            result = self._run_brahma_multiaction(effective_text)
            summary = str(result or "").strip() or "Done."
            self.write_log(f"Brahma AI: {summary}")
            self._speak_response_async(summary)
            return {"ok": True, "summary": summary}
        except Exception as exc:
            summary = str(exc) or "Advanced command failed."
            self.write_log(f"[error] {summary}")
            return {"ok": False, "summary": summary}
        finally:
            self.advanced_terminal_active = False
            self.status_text = "ONLINE" if not self.mic_enabled else "LISTENING"

    def get_connection_info(self):
        addresses = []

        def add_ip(ip):
            if not ip or ip == "0.0.0.0":
                return
            try:
                parsed = ipaddress.ip_address(ip)
            except ValueError:
                return
            if parsed.version != 4 or parsed.is_unspecified:
                return
            if ip not in addresses:
                addresses.append(ip)

        def ip_priority(ip):
            try:
                parsed = ipaddress.ip_address(ip)
            except ValueError:
                return (9, ip)

            if parsed.is_loopback:
                return (5, ip)
            if parsed.is_link_local:
                return (4, ip)
            if parsed.is_private:
                return (0, ip)
            if parsed.is_global:
                return (1, ip)
            return (3, ip)

        try:
            probe = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            probe.settimeout(0.2)
            probe.connect(("8.8.8.8", 80))
            add_ip(probe.getsockname()[0])
            probe.close()
        except Exception:
            pass

        try:
            hostnames = {socket.gethostname()}
            try:
                hostnames.add(socket.getfqdn())
            except Exception:
                pass
            for name in hostnames:
                for result in socket.getaddrinfo(name, None, socket.AF_INET, socket.SOCK_STREAM):
                    add_ip(result[4][0])
        except Exception:
            pass

        addresses.sort(key=ip_priority)
        primary_ip = addresses[0] if addresses else "127.0.0.1"
        host_urls = [f"http://{ip}:{PORT}" for ip in addresses] or [f"http://127.0.0.1:{PORT}"]
        primary_host = f"http://{primary_ip}:{PORT}"
        return {
            "ok": True,
            "port": PORT,
            "listenHost": HOST,
            "primaryHost": primary_host,
            "qrUrl": f"brahma://connect?host={primary_host}",
            "addresses": host_urls,
        }

    def create_project_workspace(self, name: str):
        safe_name = str(name or "").strip()
        if not safe_name:
            return {"ok": False, "error": "missing_name"}
        safe_name = "".join("_" if ch in '<>:"/\\|?*' or ord(ch) < 32 else ch for ch in safe_name).strip()
        if not safe_name:
            return {"ok": False, "error": "invalid_name"}
        root = Path.home() / "Documents" / "Brahma Projects"
        target = root / safe_name
        try:
            target.mkdir(parents=True, exist_ok=True)
        except Exception as exc:
            return {"ok": False, "error": str(exc)}
        self.save_hybrid_settings({
            "project_workspace_path": str(target),
            "project_workspace_name": target.name,
        })
        return {"ok": True, "path": str(target), "name": target.name}

    def set_project_workspace(self, path_value: str):
        target = Path(str(path_value or "").strip()).expanduser()
        if not str(target):
            return {"ok": False, "error": "missing_path"}
        if not target.exists() or not target.is_dir():
            return {"ok": False, "error": "folder_not_found"}
        self.save_hybrid_settings({
            "project_workspace_path": str(target),
            "project_workspace_name": target.name,
        })
        return {"ok": True, "path": str(target), "name": target.name}

    def open_project_workspace(self):
        target = Path(str(self.hybrid_settings.get("project_workspace_path") or "").strip()).expanduser()
        if not target.exists():
            return {"ok": False, "error": "folder_not_found"}
        try:
            os.startfile(str(target))
        except Exception as exc:
            return {"ok": False, "error": str(exc)}
        return {"ok": True, "path": str(target), "name": target.name}

    def _api_keys_exist(self):
        return API_FILE.exists() or ACTION_API_FILE.exists() or LEGACY_ACTION_FILE.exists()

    def _ensure_primary_api_file(self):
        """Copy an existing API key into the current config dir if missing."""
        if API_FILE.exists():
            return
        for candidate in (ACTION_API_FILE, LEGACY_ACTION_FILE):
            if candidate.exists():
                try:
                    API_FILE.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy(candidate, API_FILE)
                    return
                except Exception:
                    continue

    def _prime_api_key_env(self):
        for key_path in (API_FILE, ACTION_API_FILE, LEGACY_ACTION_FILE):
            if key_path.exists():
                try:
                    key = json.load(open(key_path, "r", encoding="utf-8")).get("gemini_api_key", "").strip()
                    if key:
                        os.environ["GOOGLE_API_KEY"] = key
                        os.environ["GEMINI_API_KEY"] = key
                        return
                except Exception:
                    continue

    def wait_for_api_key(self):
        while not self._api_key_ready:
            time.sleep(0.1)

    def set_send_callback(self, cb):
        self.send_callback = cb

    def write_log(self, text: str, mirror: bool = True):
        text = str(text or "")
        if self.advanced_terminal_active and any(
            marker in text
            for marker in (
                "Live session error:",
                "timed out during opening handshake",
                "getaddrinfo failed",
                "unhandled errors in a TaskGroup",
            )
        ):
            return
        with self._lock:
            self.logs.append({
                "ts": time.time(),
                "text": text,
            })
        if LOG_FILE:
            try:
                Path(LOG_FILE).parent.mkdir(parents=True, exist_ok=True)
                with open(LOG_FILE, "a", encoding="utf-8") as f:
                    f.write(f"{time.strftime('%Y-%m-%d %H:%M:%S')} {text}\n")
            except Exception:
                pass
        lower = text.lower()
        if lower.startswith("[error]"):
            self.status_text = "ERROR"
        elif lower.startswith("brahma:") or lower.startswith("[ai]"):
            self.status_text = "RESPONDING"
        elif lower.startswith("you:"):
            self.status_text = "PROCESSING"

        # Speak command/assistant/system responses exactly once.
        speak_text_line = ""
        if text.startswith("Brahma AI:"):
            speak_text_line = text.split("Brahma AI:", 1)[1].strip()
        elif text.startswith("You:"):
            speak_text_line = text.split("You:", 1)[1].strip()
        elif text.startswith("SYS:"):
            speak_text_line = text.split("SYS:", 1)[1].strip()
        elif text.startswith("[sys]"):
            speak_text_line = text.split("[sys]", 1)[1].strip()
        elif text.startswith("[browser]"):
            speak_text_line = text.split("[browser]", 1)[1].strip()
        elif text.startswith("[error]"):
            speak_text_line = text.split("[error]", 1)[1].strip()
        elif text.startswith("[adv]"):
            speak_text_line = text.split("[adv]", 1)[1].strip()

        if speak_text_line:
            now = time.time()
            if (
                speak_text_line
                and (speak_text_line != self._last_spoken_message or (now - self._last_spoken_ts) > 1.0)
            ):
                self._last_spoken_message = speak_text_line
                self._last_spoken_ts = now
                self._speak_response_async(speak_text_line)

        if mirror:
            try:
                self._mirror_discord_message(text)
            except Exception:
                pass

    def start_speaking(self):
        self.status_text = "SPEAKING"

    def stop_speaking(self):
        self.status_text = "ONLINE"

    def _mirror_discord_message(self, text: str):
        text = str(text or "").strip()
        if not text:
            return
        settings = dict(self.discord_settings or {})
        if not settings.get("mirror_enabled"):
            return
        channel_ids = [
            str(item).strip()
            for item in settings.get("remote_channel_ids") or []
            if str(item).strip()
        ]
        if not channel_ids:
            return
        if not self._discord_token():
            return
        if not (text.startswith("You:") or text.startswith("Brahma AI:")):
            return
        for channel_id in channel_ids:
            try:
                self._discord_send_channel_message(channel_id, text)
            except Exception:
                continue

    def process_command_text(self, text: str, source: str = "local", discord_channel_id: str | None = None):
        text = str(text or "").strip()
        if not text:
            return ""
        if source == "discord":
            self._discord_last_command = text
            self._discord_last_command_ts = time.time()
        mirror = source != "discord"
        self.write_log(f"You: {text}", mirror=mirror)
        direct_result = self.try_handle_direct_command(text)
        if direct_result:
            response = str(direct_result)
            self.write_log(f"Brahma AI: {response}", mirror=mirror)
            return response
        plugin_response, plugin_meta = self.plugin_manager.handle(text, {"source": source})
        if plugin_response:
            response = str(plugin_response)
            if plugin_meta:
                self.write_log(f"[plugin:{plugin_meta.get('name')}] {response}", mirror=mirror)
            else:
                self.write_log(f"[plugin] {response}", mirror=mirror)
            return response
        route = route_command(text, advanced_mode=False)

        if self.send_callback and self.live_ready:
            self.send_callback(self._inject_project_context(text))
            return "Command received. Executing..."

        result, _route = self.execute_hybrid_command(text, force_openclaw=False)
        response = str(result or "Done.")
        self.write_log(f"Brahma AI: {response}", mirror=mirror)
        return response

    def set_live_user_text(self, text: str):
        with self._lock:
            self.live_user_text = text or ""
        if self.live_user_text:
            self.status_text = "PROCESSING"

    def set_live_ai_text(self, text: str):
        with self._lock:
            self.live_ai_text = text or ""
        if self.live_ai_text:
            self.status_text = "SPEAKING"

    def clear_live_transcripts(self):
        with self._lock:
            self.live_user_text = ""
            self.live_ai_text = ""
        if self.mic_enabled:
            self.status_text = "LISTENING"
        else:
            self.status_text = "ONLINE"

    def start_screen_analysis(self):
        self.screen_analysis_active = True
        self.status_text = "ANALYZING"

    def stop_screen_analysis(self):
        self.screen_analysis_active = False
        self.status_text = "ONLINE"

    def capture_screen(self):
        """Capture primary monitor and cache base64 preview for UI."""
        try:
            import mss
            import mss.tools
        except Exception:
            return None
        try:
            with mss.mss() as sct:
                monitor = sct.monitors[1] if len(sct.monitors) > 1 else sct.monitors[0]
                shot = sct.grab(monitor)
                img_bytes = mss.tools.to_png(shot.rgb, shot.size)
                b64 = base64.b64encode(img_bytes).decode("ascii")
                with self._lock:
                    self.last_screenshot_b64 = b64
                return b64
        except Exception:
            return None

    def show_gesture_cursor(self, *args, **kwargs):
        return

    def move_gesture_cursor(self, *args, **kwargs):
        return

    def flash_gesture_cursor(self, *args, **kwargs):
        return

    def hide_gesture_cursor(self, *args, **kwargs):
        return

    def update_gesture_hand(self, *args, **kwargs):
        return

    def update_camera_preview_frame(self, *args, **kwargs):
        if not args:
            return
        frame_rgb = args[0]
        try:
            import cv2
            frame_bgr = cv2.cvtColor(frame_rgb, cv2.COLOR_RGB2BGR)
            ok, buf = cv2.imencode(".jpg", frame_bgr, [int(cv2.IMWRITE_JPEG_QUALITY), 55])
            if ok:
                with self._lock:
                    self.camera_preview_b64 = base64.b64encode(buf.tobytes()).decode("ascii")
        except Exception:
            return

    def gesture_close_app(self):
        try:
            if sys.platform == "win32":
                subprocess.Popen(["taskkill", "/F", "/IM", "electron.exe"], shell=False)
                return True
        except Exception:
            pass
        return False

    def save_api_key(self, key: str):
        key = key.strip()
        if not key:
            return False
        # Primary location (AppData or configured BRAHMA_CONFIG_DIR)
        API_FILE.parent.mkdir(parents=True, exist_ok=True)
        with open(API_FILE, "w", encoding="utf-8") as f:
            json.dump({"gemini_api_key": key}, f, indent=4)
        # Additional copies for packaged/local modules
        try:
            ACTION_API_FILE.parent.mkdir(parents=True, exist_ok=True)
            with open(ACTION_API_FILE, "w", encoding="utf-8") as f:
                json.dump({"gemini_api_key": key}, f, indent=4)
        except Exception:
            pass
        try:
            LEGACY_ACTION_FILE.parent.mkdir(parents=True, exist_ok=True)
            with open(LEGACY_ACTION_FILE, "w", encoding="utf-8") as f:
                json.dump({"gemini_api_key": key}, f, indent=4)
        except Exception:
            pass
        self._api_key_ready = True
        os.environ["GOOGLE_API_KEY"] = key
        self.write_log("SYS: API key saved.")
        return True

    def toggle_gesture(self):
        if not self.gesture_enabled:
            result = self.gesture_controller.start()
            if "enabled" in result.lower():
                self.gesture_enabled = True
            return result
        result = self.gesture_controller.stop()
        self.gesture_enabled = False
        return result

    def toggle_mic(self):
        self.mic_enabled = not self.mic_enabled
        if self.mic_enabled:
            self.write_log("[sys] Microphone enabled.")
            self.status_text = "LISTENING"
            return "Microphone enabled."
        self.write_log("[sys] Microphone disabled.")
        if self.status_text == "LISTENING":
            self.status_text = "ONLINE"
        return "Microphone disabled."

    def is_mic_enabled(self):
        return self.mic_enabled

    def get_state(self):
        with self._lock:
            logs = list(self.logs)
        return {
            "status": self.status_text,
            "apiKeyReady": self._api_key_ready,
            "screenAnalysisActive": self.screen_analysis_active,
            "gestureEnabled": self.gesture_enabled,
            "micEnabled": self.mic_enabled,
            "cameraPreview": self.camera_preview_b64,
            "liveUserText": self.live_user_text,
            "liveAiText": self.live_ai_text,
            "savedSequences": self.saved_sequences,
            "routineSchedules": self.sequence_schedules,
            "kasaDevices": self.kasa_devices,
            "voiceSettings": self.voice_settings,
            "voiceCapabilities": detect_voice_capabilities(self.voice_settings),
            "logs": logs,
            "hybrid": self.get_hybrid_state(),
            "gmailConfigured": bool(self.gmail_settings.get("address") and self.gmail_settings.get("app_password")),
            "discordConfigured": bool(self.discord_settings.get("bot_token")),
            "discordSettings": {
                "remoteEnabled": bool(self.discord_settings.get("remote_enabled")),
                "mirrorEnabled": bool(self.discord_settings.get("mirror_enabled", True)),
                "remoteChannelIds": list(self.discord_settings.get("remote_channel_ids") or []),
            },
            "discordRemote": {
                "lastCommand": self._discord_last_command,
                "lastCommandTs": self._discord_last_command_ts,
            },
            "plugins": self.plugin_manager.list_plugins(),
            "automation": {
                "mode": self.automation_mode,
                "plan": self.automation_plan,
                "pendingConfirmation": self.pending_confirmation,
                "lastResult": self.last_automation_result,
                "lastScreenshot": self.last_screenshot_b64,
            },
            "advancedTerminalActive": self.advanced_terminal_active,
        }

    def save_voice_settings(self, next_settings):
        self.voice_settings = save_voice_settings(self._edge_only_voice_settings(next_settings))
        return self.voice_settings

    def test_voice(self, text="Brahma AI voice systems online. How can I help?"):
        result = speak_text(text, settings=self._edge_only_voice_settings(self.voice_settings), player=self)
        return result

    def refresh_sequences(self):
        result = workflow_manager({"action": "list"}, player=self)
        if result.startswith("Saved workflows:"):
            names = result.split(":", 1)[1].strip()
            self.saved_sequences = [item.strip() for item in names.split(",") if item.strip()]
        else:
            self.saved_sequences = []

    def refresh_kasa_devices(self):
        result = kasa_control({"action": "discover"}, player=self)
        if result.startswith("Kasa devices found:"):
            lines = result.splitlines()[1:]
            self.kasa_devices = [line.split(". ", 1)[1] for line in lines if ". " in line]
        else:
            self.kasa_devices = []


UI = BridgeUI()


def _extract_json_block(text: str):
    """Best-effort extract JSON object/array from LLM response."""
    if not text:
        return None
    start = text.find("{")
    if start == -1:
        return None
    end = text.rfind("}")
    if end == -1:
        return None
    snippet = text[start : end + 1]
    try:
        return json.loads(snippet)
    except Exception:
        return None


def _normalize_bbox(bbox, screen_width, screen_height):
    if not bbox:
        return None
    try:
        x = float(bbox.get("x", 0))
        y = float(bbox.get("y", 0))
        w = float(bbox.get("w", 0) or bbox.get("width", 0))
        h = float(bbox.get("h", 0) or bbox.get("height", 0))
    except Exception:
        return None
    # If values look normalized (<1), scale to pixels
    if x <= 1 and y <= 1 and (w <= 1 or h <= 1):
        x *= screen_width
        y *= screen_height
        w *= screen_width
        h *= screen_height
    return {
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "cx": x + w / 2,
        "cy": y + h / 2,
    }


def plan_actions_from_screen(task_text: str, screenshot_b64: str):
    """Use Gemini Vision (if available) to propose UI actions."""
    plan_id = f"plan-{int(time.time() * 1000)}"
    if not screenshot_b64:
        return {
            "id": plan_id,
            "task": task_text,
            "actions": [],
            "note": "No screenshot captured; vision planner skipped.",
        }
    try:
        from google import genai

        api_key = os.environ.get("GOOGLE_API_KEY", "").strip()
        if not api_key:
            return {
                "id": plan_id,
                "task": task_text,
                "actions": [],
                "note": "Gemini API key missing. Add key in settings.",
            }

        client = genai.Client(api_key=api_key, http_options={"api_version": "v1"})
        prompt = (
            "You are a desktop automation planner. Given a screenshot and a user command, "
            "return a JSON object with an 'actions' array. Each action has: "
            "{id, kind: one of [click, type, scroll], text (for type), "
            "bbox: {x,y,w,h} in normalized [0,1] coordinates of the visible screen, "
            "note (optional)}. Return ONLY JSON."
        )
        image_bytes = base64.b64decode(screenshot_b64)
        resp = client.models.generate_content(
            model="models/gemini-2.5-flash",
            contents=[
                {"role": "user", "parts": [{"text": prompt}]},
                {"role": "user", "parts": [{"text": f"Task: {task_text}"}, {"inline_data": {"mime_type": "image/png", "data": image_bytes}}]},
            ],
        )
        text = resp.text or ""
        parsed = _extract_json_block(text)
        if not parsed:
            return {
                "id": plan_id,
                "task": task_text,
                "actions": [],
                "note": "Vision response not in JSON; skipping actions.",
                "raw": text,
            }
        parsed["id"] = parsed.get("id") or plan_id
        parsed["task"] = task_text
        return parsed
    except Exception:
        _log_traceback("[error] Vision planning failed.")
        return {
            "id": plan_id,
            "task": task_text,
            "actions": [],
            "note": "Vision planning error; see logs.",
        }


def execute_actions(plan: dict, allow_ids=None):
    """Execute planned actions with pyautogui, respecting an allow list."""
    allow_set = set(allow_ids) if allow_ids else None
    actions = plan.get("actions", []) if plan else []
    results = []
    try:
        import pyautogui
    except Exception:
        return [{"id": a.get("id"), "ok": False, "error": "pyautogui not available"} for a in actions]

    screen_w, screen_h = pyautogui.size()
    for action in actions:
        action_id = action.get("id") or str(uuid.uuid4())
        if allow_set is not None and action_id not in allow_set:
            continue
        kind = (action.get("kind") or "").lower()
        try:
            if kind == "click":
                bbox = _normalize_bbox(action.get("bbox"), screen_w, screen_h)
                if not bbox:
                    results.append({"id": action_id, "ok": False, "error": "bbox missing"})
                    continue
                pyautogui.moveTo(bbox["cx"], bbox["cy"], duration=0.15)
                pyautogui.click()
                results.append({"id": action_id, "ok": True, "kind": "click"})
            elif kind == "type":
                text = action.get("text", "")
                pyautogui.write(text, interval=0.02)
                results.append({"id": action_id, "ok": True, "kind": "type"})
            elif kind == "scroll":
                amount = int(action.get("amount", -600))
                pyautogui.scroll(amount)
                results.append({"id": action_id, "ok": True, "kind": "scroll"})
            else:
                results.append({"id": action_id, "ok": False, "error": f"unknown kind {kind}"})
            time.sleep(0.12)
        except Exception as exc:
            results.append({"id": action_id, "ok": False, "error": str(exc)})
    return results


def _log_traceback(prefix: str):
    """Write a traceback to the configured log file and stdout for quick diagnosis."""
    try:
        UI.write_log(prefix)
    except Exception:
        pass
    tb = traceback.format_exc()
    if LOG_FILE:
        try:
            Path(LOG_FILE).parent.mkdir(parents=True, exist_ok=True)
            with open(LOG_FILE, "a", encoding="utf-8") as f:
                f.write(tb + "\n")
        except Exception:
            pass
    print(tb)


class ApiHandler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        # Silence default stdout/stderr logging to avoid invalid handle errors in packaged/hidden contexts.
        return

    def _send_cors_headers(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type, Authorization, X-Requested-With")
        self.send_header("Access-Control-Allow-Private-Network", "true")

    def _send_json(self, data, status=200):
        body = json.dumps(data).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self._send_cors_headers()
        self.end_headers()
        self.wfile.write(body)

    def do_OPTIONS(self):
        self.send_response(204)
        self._send_cors_headers()
        self.end_headers()

    def do_GET(self):
        try:
            if self.path == "/api/state":
                self._send_json(UI.get_state())
                return
            if self.path == "/api/connection-info":
                self._send_json(UI.get_connection_info())
                return
            if self.path == "/api/automation/state":
                self._send_json({
                    "ok": True,
                    "automation": UI.get_state().get("automation", {}),
                })
                return
            if self.path == "/api/plugins":
                self._send_json({
                    "ok": True,
                    "plugins": UI.plugin_manager.list_plugins(),
                })
                return
            if self.path == "/health":
                self._send_json({"ok": True, **UI.get_connection_info()})
                return
            self._send_json({"error": "not_found"}, 404)
        except Exception:
            _log_traceback("[error] GET handler failed")
            try:
                self._send_json({"error": "server_error"}, 500)
            except Exception:
                pass

    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", "0"))
            raw = self.rfile.read(length) if length else b"{}"
            try:
                data = json.loads(raw.decode("utf-8"))
            except Exception:
                data = {}

        except Exception:
            _log_traceback("[error] Failed to parse request body")
            data = {}

        try:
            if self.path == "/api/send":
                text = (data.get("text") or "").strip()
                if not text:
                    self._send_json({"ok": False, "error": "missing_text"}, 400)
                    return
                incoming_project_path = str(data.get("projectWorkspacePath") or "").strip()
                incoming_project_name = str(data.get("projectWorkspaceName") or "").strip()
                if incoming_project_path != str(UI.hybrid_settings.get("project_workspace_path") or "").strip() or incoming_project_name != str(UI.hybrid_settings.get("project_workspace_name") or "").strip():
                    UI.save_hybrid_settings({
                        "project_workspace_path": incoming_project_path,
                        "project_workspace_name": incoming_project_name,
                    })
                UI.write_log(f"You: {text}")
                direct_result = UI.try_handle_direct_command(text)
                if direct_result:
                    UI.write_log(f"Brahma AI: {direct_result}")
                    self._send_json({
                        "ok": True,
                        "mode": "direct",
                        "message": direct_result,
                        "engine": "direct_handler",
                    })
                    return
                plugin_response, plugin_meta = UI.plugin_manager.handle(text, {"source": "ui"})
                if plugin_response:
                    message = str(plugin_response)
                    if plugin_meta:
                        UI.write_log(f"[plugin:{plugin_meta.get('name')}] {message}")
                    else:
                        UI.write_log(f"[plugin] {message}")
                    self._send_json({
                        "ok": True,
                        "mode": "plugin",
                        "message": message,
                        "engine": "plugin",
                    })
                    return
                route = route_command(text, advanced_mode=False)
                if UI.send_callback and UI.live_ready:
                    UI.send_callback(UI._inject_project_context(text))
                    self._send_json({
                        "ok": True,
                        "mode": "online",
                        "engine": "brahma_live",
                        "route": {
                            "engine": route.engine,
                            "confidence": route.confidence,
                            "reason": route.reason,
                        },
                    })
                    return

                result, route = UI.execute_hybrid_command(text, force_openclaw=False)
                UI.write_log(f"Brahma AI: {result}")
                self._send_json({
                    "ok": True,
                    "mode": "offline",
                    "message": result,
                    "engine": route.get("engine"),
                    "route": route,
                })
                return

            if self.path == "/api/api-key":
                ok = UI.save_api_key(data.get("key", ""))
                self._send_json({"ok": ok}, 200 if ok else 400)
                return

            if self.path == "/api/toggle-gesture":
                result = UI.toggle_gesture()
                self._send_json({"ok": True, "message": result, "gestureEnabled": UI.gesture_enabled})
                return

            if self.path == "/api/toggle-mic":
                result = UI.toggle_mic()
                self._send_json({"ok": True, "message": result, "micEnabled": UI.mic_enabled})
                return

            if self.path == "/api/discord-settings":
                token = str(data.get("botToken") if "botToken" in data else data.get("bot_token", "")).strip()
                if "botToken" in data or "bot_token" in data:
                    UI.discord_settings["bot_token"] = token
                if "remoteEnabled" in data:
                    UI.discord_settings["remote_enabled"] = bool(data.get("remoteEnabled"))
                if "mirrorEnabled" in data:
                    UI.discord_settings["mirror_enabled"] = bool(data.get("mirrorEnabled"))
                channel_ids = data.get("remoteChannelIds", UI.discord_settings.get("remote_channel_ids", []))
                if isinstance(channel_ids, str):
                    channel_ids = [item.strip() for item in channel_ids.split(",") if item.strip()]
                if isinstance(channel_ids, list):
                    UI.discord_settings["remote_channel_ids"] = [str(item).strip() for item in channel_ids if str(item).strip()]
                UI._save_discord_settings()
                self._send_json({
                    "ok": True,
                    "discordConfigured": bool(UI.discord_settings.get("bot_token")),
                    "discordSettings": {
                        "remoteEnabled": bool(UI.discord_settings.get("remote_enabled")),
                        "mirrorEnabled": bool(UI.discord_settings.get("mirror_enabled", True)),
                        "remoteChannelIds": list(UI.discord_settings.get("remote_channel_ids") or []),
                    },
                })
                return

            if self.path == "/api/plugins/reload":
                UI.plugin_manager.reload()
                self._send_json({
                    "ok": True,
                    "plugins": UI.plugin_manager.list_plugins(),
                })
                return

            if self.path == "/api/plugins/open-folder":
                try:
                    PLUGINS_DIR.mkdir(parents=True, exist_ok=True)
                    if sys.platform == "win32":
                        os.startfile(str(PLUGINS_DIR))  # type: ignore[attr-defined]
                    else:
                        subprocess.Popen(["xdg-open", str(PLUGINS_DIR)])
                    self._send_json({"ok": True})
                except Exception as exc:
                    self._send_json({"ok": False, "error": str(exc)}, 500)
                return

            if self.path == "/api/discord-test":
                start = time.time()
                bot = UI._discord_get_bot_user()
                latency_ms = int((time.time() - start) * 1000)
                channel_name = ""
                channel_id = ""
                channel_ids = UI.discord_settings.get("remote_channel_ids") or []
                if channel_ids:
                    channel_id = str(channel_ids[0])
                    try:
                        channel_info = UI._discord_api_request("GET", f"/channels/{channel_id}")
                        channel_name = str(channel_info.get("name") or "")
                    except Exception:
                        channel_name = ""
                self._send_json({
                    "ok": True,
                    "latencyMs": latency_ms,
                    "botName": str((bot or {}).get("username") or ""),
                    "channelId": channel_id,
                    "channelName": channel_name,
                })
                return

            if self.path == "/api/voice-settings":
                UI.save_voice_settings(data)
                self._send_json({
                    "ok": True,
                    "voiceSettings": UI.voice_settings,
                    "voiceCapabilities": detect_voice_capabilities(UI.voice_settings),
                })
                return

            if self.path == "/api/hybrid-settings":
                settings = UI.save_hybrid_settings({
                    "project_workspace_path": data.get("projectWorkspacePath"),
                    "project_workspace_name": data.get("projectWorkspaceName"),
                })
                self._send_json({
                    "ok": True,
                    "hybrid": UI.get_hybrid_state(),
                    "settings": settings,
                })
                return

            if self.path == "/api/project-workspace":
                action = str(data.get("action") or "").strip().lower()
                if action == "create":
                    result = UI.create_project_workspace(data.get("name"))
                elif action == "select":
                    result = UI.set_project_workspace(data.get("path"))
                elif action == "open":
                    result = UI.open_project_workspace()
                elif action == "clear":
                    UI.save_hybrid_settings({
                        "project_workspace_path": "",
                        "project_workspace_name": "",
                    })
                    result = {"ok": True}
                else:
                    result = {"ok": False, "error": "unsupported_action"}
                status = 200 if result.get("ok") else 400
                self._send_json({
                    "ok": bool(result.get("ok")),
                    "result": result,
                    "hybrid": UI.get_hybrid_state(),
                    "error": result.get("error"),
                }, status=status)
                return

            if self.path == "/api/test-voice":
                result = UI.test_voice(str(data.get("text") or "").strip())
                status = 200 if result.get("ok") else 400
                self._send_json(result, status=status)
                return

            if self.path == "/api/kasa":
                action = str(data.get("action") or "discover").strip().lower()
                device_name = str(data.get("device_name") or "").strip()
                result = kasa_control({"action": action, "device_name": device_name}, player=UI)
                UI.refresh_kasa_devices()
                UI.write_log(f"Brahma AI: {result}")
                self._send_json({"ok": True, "message": result, "kasaDevices": UI.kasa_devices})
                return

            if self.path == "/api/automation/mode":
                mode = str(data.get("mode") or "").strip().lower() or "assist"
                if mode not in ("observe", "assist", "do"):
                    mode = "assist"
                UI.automation_mode = mode
                self._send_json({"ok": True, "mode": UI.automation_mode})
                return

            if self.path == "/api/automation/intent":
                text = (data.get("text") or "").strip()
                mode = str(data.get("mode") or UI.automation_mode).strip().lower() or "assist"
                if mode not in ("observe", "assist", "do"):
                    mode = "assist"
                UI.automation_mode = mode
                if not text:
                    self._send_json({"ok": False, "error": "missing_text"}, 400)
                    return
                UI.write_log(f"[auto] {mode.upper()} request: {text}")
                UI.start_screen_analysis()
                screenshot = UI.capture_screen()
                plan = plan_actions_from_screen(text, screenshot)
                UI.stop_screen_analysis()
                UI.automation_plan = plan
                UI.pending_confirmation = bool(plan.get("actions"))
                UI.last_automation_result = None
                self._send_json({
                    "ok": True,
                    "plan": plan,
                    "requireConfirm": UI.pending_confirmation,
                    "screenshot": bool(screenshot),
                })
                return

            if self.path == "/api/automation/confirm":
                plan_id = str(data.get("planId") or "")
                if not UI.automation_plan or plan_id != UI.automation_plan.get("id"):
                    self._send_json({"ok": False, "error": "plan_mismatch"}, 400)
                    return
                allow_ids = data.get("allow") or None
                results = execute_actions(UI.automation_plan, allow_ids)
                UI.pending_confirmation = False
                UI.last_automation_result = results
                for res in results:
                    status = "ok" if res.get("ok") else f"fail ({res.get('error')})"
                    UI.write_log(f"[auto] action {res.get('id')}: {status}")
                self._send_json({"ok": True, "results": results})
                return

            if self.path == "/api/automation/cancel":
                UI.pending_confirmation = False
                UI.automation_plan = None
                UI.last_automation_result = None
                self._send_json({"ok": True})
                return

            if self.path == "/api/cad/generate":
                prompt = (data.get("prompt") or "").strip()
                if not prompt:
                    self._send_json({"ok": False, "error": "missing_prompt"}, 400)
                    return
                UI.write_log("[cad] CAD generation is currently unavailable (engine removed).")
                self._send_json({
                    "ok": False,
                    "error": "cad_unavailable",
                }, 503)
                return

            if self.path == "/api/sequence":
                action = str(data.get("action") or "").strip().lower()
                name = str(data.get("name") or "").strip()
                steps = data.get("steps") or []
                auto_run = bool(data.get("autoRun", False))
                schedule_time = str(data.get("scheduleTime") or data.get("runAt") or "").strip()

                if action == "save":
                    if auto_run and not re.fullmatch(r"\d{2}:\d{2}", schedule_time):
                        self._send_json({"ok": False, "error": "invalid_schedule_time"}, 400)
                        return
                    result = workflow_manager(
                        {"action": "create", "name": name, "steps": steps},
                        player=UI,
                    )
                    if name:
                        if auto_run and re.fullmatch(r"\d{2}:\d{2}", schedule_time):
                            UI._set_sequence_schedule(name, schedule_time, True)
                        elif "autoRun" in data and not auto_run:
                            UI._delete_sequence_schedule(name)
                    UI.refresh_sequences()
                    UI.write_log(f"Brahma AI: {result}")
                    self._send_json(
                        {
                            "ok": True,
                            "message": result,
                            "savedSequences": UI.saved_sequences,
                            "routineSchedules": UI.sequence_schedules,
                        }
                    )
                    return

                if action == "start":
                    step_list = [str(item).strip() for item in steps] if isinstance(steps, list) else []
                    step_list = [item for item in step_list if item]
                    if name and step_list:
                        workflow_manager(
                            {"action": "create", "name": name, "steps": step_list},
                            player=UI,
                        )
                    result = workflow_manager(
                        {"action": "run", "name": name, "step_runner": UI._execute_routine_step},
                        player=UI,
                    )
                    UI.write_log(f"Brahma AI: {result}")
                    self._send_json({"ok": True, "message": result, "routineSchedules": UI.sequence_schedules})
                    return

                if action == "delete":
                    result = workflow_manager(
                        {"action": "delete", "name": name},
                        player=UI,
                    )
                    if name:
                        UI._delete_sequence_schedule(name)
                    UI.refresh_sequences()
                    UI.write_log(f"Brahma AI: {result}")
                    self._send_json(
                        {
                            "ok": True,
                            "message": result,
                            "savedSequences": UI.saved_sequences,
                            "routineSchedules": UI.sequence_schedules,
                        }
                    )
                    return

                if action == "list":
                    UI.refresh_sequences()
                    self._send_json(
                        {
                            "ok": True,
                            "savedSequences": UI.saved_sequences,
                            "routineSchedules": UI.sequence_schedules,
                        }
                    )
                    return

            self._send_json({"error": "not_found"}, 404)
        except Exception:
            _log_traceback("[error] POST handler failed")
            try:
                self._send_json({"error": "server_error"}, 500)
            except Exception:
                pass


def run_backend():
    try:
        UI.wait_for_api_key()
        UI.refresh_sequences()
        UI.refresh_kasa_devices()
        brahma = JarvisLive(UI)
        UI.set_send_callback(brahma.submit_text)
        UI.write_log("SYS: Backend event loop starting.")
        asyncio.run(brahma.run())
    except Exception:
        _log_traceback("[error] Backend crashed; see stack trace below.")


def main():
    threading.Thread(target=run_backend, daemon=True, name="BrahmaBackend").start()
    server = ThreadingHTTPServer((HOST, PORT), ApiHandler)
    try:
        print(f"[Bridge] Electron backend listening on http://{HOST}:{PORT}")
    except Exception:
        pass
    server.serve_forever()


if __name__ == "__main__":
    main()
